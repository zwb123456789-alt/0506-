# -*- coding: utf-8 -*-
"""生成 0603 项目进展汇报 PPT（绿色学术风，16:9，含每页备注）。"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

ROOT = r"d:\我的文件\研究生学术\光学项目\0506新"
GDIR = os.path.join(ROOT, "汇报0603_Claude生成PPT指导")
OUT = os.path.join(ROOT, "论文改进", "汇报材料", "20260603_项目进展汇报.pptx")

# ---- 配色（取自模板 theme1.xml）----
GREEN  = RGBColor(0x01, 0x6A, 0x3F)   # 主深绿
LGREEN = RGBColor(0x75, 0xBD, 0x42)   # 亮绿
TEAL   = RGBColor(0x30, 0xC0, 0xB4)   # 青
PALE   = RGBColor(0xE3, 0xF1, 0xD9)   # 极浅绿（卡片底）
INK    = RGBColor(0x26, 0x2B, 0x2E)   # 近黑正文
GRAY   = RGBColor(0x6B, 0x72, 0x70)   # 次要灰
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
RED    = RGBColor(0xC0, 0x39, 0x2B)   # 风险/崩溃
CREAM  = RGBColor(0xF6, 0xF8, 0xF4)   # 浅底

HEAD = "Microsoft YaHei"   # 标题
BODY = "Microsoft YaHei"   # 正文

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def _set_font(run, size, color=INK, bold=False, font=BODY, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color
    # 东亚字体
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font)


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """lines: list of (text, size, color, bold, font, italic, space_after)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    for i, ln in enumerate(lines):
        text, size, color, bold = ln[0], ln[1], ln[2], ln[3]
        font = ln[4] if len(ln) > 4 else BODY
        italic = ln[5] if len(ln) > 5 else False
        sa = ln[6] if len(ln) > 6 else 4
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sa)
        p.space_before = Pt(0)
        # 支持单段多 run：text 可为 list[(t,bold,color)]
        if isinstance(text, list):
            for seg in text:
                r = p.add_run(); r.text = seg[0]
                _set_font(r, size, seg[2] if len(seg) > 2 else color,
                          seg[1] if len(seg) > 1 else bold, font, italic)
        else:
            r = p.add_run(); r.text = text
            _set_font(r, size, color, bold, font, italic)
    return tb


def bullets(slide, x, y, w, h, items, size=15, color=INK, gap=7, marker=True, mcolor=None):
    """items: list of str 或 (text, bold). 以绿色方块作 marker。"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1)
    mcolor = mcolor or LGREEN
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0)
        p.alignment = PP_ALIGN.LEFT
        if marker:
            rb = p.add_run(); rb.text = "▍ "
            _set_font(rb, size, mcolor, True, BODY)
        if isinstance(it, tuple):
            r = p.add_run(); r.text = it[0]
            _set_font(r, size, color, it[1] if len(it) > 1 else False, BODY)
        elif isinstance(it, list):
            for seg in it:
                r = p.add_run(); r.text = seg[0]
                _set_font(r, size, seg[2] if len(seg) > 2 else color,
                          seg[1] if len(seg) > 1 else False, BODY)
        else:
            r = p.add_run(); r.text = it
            _set_font(r, size, color, False, BODY)
    return tb


def pic_fit(slide, path, x, y, max_w, max_h, align="center", valign="middle"):
    """按比例放入框，返回实际 picture。"""
    im = Image.open(path); iw, ih = im.size
    ratio = iw / ih
    box_ratio = max_w / max_h
    if ratio > box_ratio:
        w = max_w; h = int(max_w / ratio)
    else:
        h = max_h; w = int(max_h * ratio)
    if align == "center":
        px = x + (max_w - w) // 2
    elif align == "left":
        px = x
    else:
        px = x + (max_w - w)
    if valign == "middle":
        py = y + (max_h - h) // 2
    elif valign == "top":
        py = y
    else:
        py = y + (max_h - h)
    return slide.shapes.add_picture(path, px, py, w, h)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def header(slide, kicker, title, page_no):
    """统一页眉：左侧色条 + kicker + 标题 + 右上页码。"""
    rect(slide, 0, 0, Inches(0.16), EMU_H, fill=GREEN)            # 左竖条
    rect(slide, Inches(0.55), Inches(0.42), Inches(0.46), Inches(0.10), fill=LGREEN)  # kicker短条
    txt(slide, Inches(0.55), Inches(0.52), Inches(9.5), Inches(0.34),
        [(kicker, 12.5, TEAL, True)], )
    txt(slide, Inches(0.55), Inches(0.80), Inches(11.8), Inches(0.72),
        [(title, 26, GREEN, True, HEAD)])
    # 页码
    txt(slide, Inches(12.35), Inches(0.46), Inches(0.7), Inches(0.34),
        [(f"{page_no:02d}", 12, GRAY, True)], align=PP_ALIGN.RIGHT)
    # 标题下细分隔（用浅底块而非下划线）
    rect(slide, Inches(0.55), Inches(1.52), Inches(12.2), Pt(1.4), fill=PALE)


P = {}  # 收集图片路径
P['fig01'] = os.path.join(ROOT, r"结果\模块C_反演\paper_summary\run_20260528_162418\fig01_bar_chart.png")
P['fig02'] = os.path.join(ROOT, r"结果\模块C_反演\paper_summary\run_20260528_162418\fig02_hit5_bar_chart.png")
P['heat'] = os.path.join(ROOT, r"结果\模块A_重构\multi_geom_ggx_yaw73_pitch37\run_20260527_195122\phase63_backscatter\fig02_ocs_heatmap.png")
P['occ'] = os.path.join(ROOT, r"结果\模块A_重构\multi_geom_ggx_yaw73_pitch37\run_20260527_195122\phase63_backscatter\fig04_occlusion_ratio_heatmap.png")
P['sat'] = os.path.join(ROOT, r"结果\模块A_重构\multi_geom_ggx_yaw73_pitch37\run_20260527_195122\phase63_backscatter\fig06_satellite_model.png")
P['resnet'] = os.path.join(GDIR, "chart_resnet_robust.png")
P['noise'] = os.path.join(GDIR, "chart_ocs_noise.png")
P['render'] = os.path.join(ROOT, r"结果\模块B_渲染\run_20260528_101944_exact_brdf\brdf_images\yaw150.00_pitch-45.00_brdf.png")

print("build helpers ready")

# ============================ 第 1 页：标题页（深绿底）============================
s = add_slide()
rect(s, 0, 0, EMU_W, EMU_H, fill=GREEN)
rect(s, 0, 0, EMU_W, Inches(0.22), fill=LGREEN)
rect(s, 0, EMU_H - Inches(0.22), EMU_W, Inches(0.22), fill=TEAL)
# 装饰：右侧大色块
rect(s, Inches(9.7), Inches(1.3), Inches(3.0), Inches(0.10), fill=TEAL)
txt(s, Inches(1.0), Inches(2.35), Inches(11.3), Inches(1.9),
    [("OCS–光度图像联合仿真与", 38, WHITE, True, HEAD, False, 6),
     ("空间目标姿态反演  项目进展汇报", 38, WHITE, True, HEAD, False, 0)])
txt(s, Inches(1.0), Inches(4.35), Inches(11.0), Inches(0.6),
    [("统一 BRDF 仿真框架 · 多模态反演实验结果 · 下一步论文推进", 18, PALE, False, BODY)])
rect(s, Inches(1.02), Inches(5.15), Inches(2.6), Pt(2), fill=LGREEN)
txt(s, Inches(1.0), Inches(5.45), Inches(11.0), Inches(0.9),
    [("汇报人：御酒的酒", 14.5, RGBColor(0xD8,0xE8,0xDC), False, BODY, False, 3),
     ("日期：2026-06-03", 14.5, RGBColor(0xD8,0xE8,0xDC), False, BODY)])
notes(s, "老师好，我这次主要汇报三个部分：第一是整个 OCS 和光度图像统一仿真框架目前做到什么程度；"
         "第二是复现和补充实验得到的主要结果；第三是论文写作现在的位置和接下来需要确认的问题。"
         "今天不展开所有代码细节，重点讲过程、结果和下一步。整个汇报大概十分钟。")

# ============================ 第 2 页：项目目标与当前主线 ============================
s = add_slide()
header(s, "PROJECT GOAL · 项目定位", "项目目标与当前主线", 2)
# 左：要点
bullets(s, Inches(0.6), Inches(1.85), Inches(6.6), Inches(4.6),
    [[("建立统一 BRDF 驱动的仿真框架", True)],
     "基于同一 STL 几何、姿态、材料、GGX BRDF 与自遮挡模型，同时生成 OCS 与光度图像",
     [("在受控仿真数据上做姿态反演基准", True)],
     "系统比较 OCS-only、image-only、Late/Feature fusion 四类方法",
     [("主线已从“某个融合模型最优”收敛为：", True)],
     [("统一物理仿真框架 ", False, GREEN), ("+ ", False, INK), ("多模态基准评估 ", False, GREEN), ("+ ", False, INK), ("条件性互补分析", False, GREEN)]],
    size=15.5, gap=9)
# 右：核心问题卡片
rect(s, Inches(7.55), Inches(1.95), Inches(5.2), Inches(4.25), fill=PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, Inches(7.55), Inches(1.95), Inches(5.2), Inches(0.62), fill=GREEN, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
txt(s, Inches(7.75), Inches(2.02), Inches(4.8), Inches(0.5),
    [("核心科学问题", 16, WHITE, True, HEAD)])
txt(s, Inches(7.85), Inches(2.85), Inches(4.65), Inches(3.2),
    [("在非均匀 BRDF、自遮挡和不同观测质量下——", 14, INK, True, BODY, False, 8),
     ([("OCS 标量光度特征", True, GREEN), (" 与 ", False, INK), ("光度图像", True, GREEN), (" 各自承载何种姿态信息？", False, INK)], 14.5, INK, False, BODY, False, 10),
     ([("多模态融合", True, TEAL), (" 在", False, INK), ("什么条件下", True, INK), ("能提供鲁棒的互补约束？", False, INK)], 14.5, INK, False, BODY, False, 0)])
notes(s, "这个项目现在的定位，已经从单纯追求某一个融合模型最优，调整为一个更稳妥的主线："
         "先建立物理一致的 OCS 和图像仿真框架，再用这个框架做受控姿态反演基准。"
         "现在的核心问题是：在非均匀 BRDF、自遮挡和不同观测质量下，OCS 和图像各自提供什么姿态信息，"
         "融合又在什么条件下真正有帮助。这个定位调整后面会解释原因。")

# ============================ 第 3 页：总体技术流程（A→B→C）============================
s = add_slide()
header(s, "PIPELINE · 总体技术流程", "三模块统一仿真与反演流程", 3)
# 顶部输入条
rect(s, Inches(0.6), Inches(1.78), Inches(12.15), Inches(0.62), fill=CREAM, line=LGREEN, line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.6), Inches(1.84), Inches(12.15), Inches(0.5),
    [([("输入　", True, GREEN), ("真实卫星 STL 三件套　·　非均匀材料分区　·　yaw–pitch 姿态网格　·　多组观测几何", False, INK)], 14, INK, False, BODY)],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 三模块卡片
cards = [
    ("模块 A", "OCS 计算", ["GGX BRDF 面元积分", "解析射线自遮挡", "5 观测几何 × 2701 姿态", "输出 OCS / 遮挡率"], GREEN),
    ("模块 B", "Blender 渲染", ["几何缓冲 MULTILAYER EXR", "Python 像素级 exact BRDF", "256 分辨率光度图像", "OPTIX GPU 0.32s/帧"], TEAL),
    ("模块 C", "姿态反演", ["OCS MLP / 图像 CNN·ResNet", "Late / Feature fusion", "mean / Hit@5° / 鲁棒性", "条件性互补分析"], LGREEN),
]
cw = Inches(3.85); gap = Inches(0.27); x0 = Inches(0.6); y0 = Inches(2.72); ch = Inches(3.05)
for i, (tag, name, items, col) in enumerate(cards):
    x = x0 + i * (cw + gap)
    rect(s, x, y0, cw, ch, fill=WHITE, line=col, line_w=1.6, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y0, cw, Inches(0.72), fill=col, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    txt(s, x + Inches(0.2), y0 + Inches(0.06), cw - Inches(0.4), Inches(0.62),
        [([(tag + "　", True, WHITE), (name, True, WHITE)], 17, WHITE, True, HEAD)], anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, x + Inches(0.22), y0 + Inches(0.92), cw - Inches(0.42), ch - Inches(1.0),
            items, size=13, gap=8, mcolor=col)
    if i < 2:
        ar = rect(s, x + cw - Inches(0.02), y0 + Inches(1.25), gap + Inches(0.06), Inches(0.5),
                  fill=col, shape=MSO_SHAPE.CHEVRON)
# 底部要点
rect(s, Inches(0.6), Inches(6.05), Inches(12.15), Inches(0.72), fill=PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.6), Inches(6.12), Inches(12.15), Inches(0.58),
    [([("关键设计：", True, GREEN), ("OCS 与图像并非两个割裂数据源，而是来自同一物理仿真框架——共享几何、姿态、材料、BRDF 与遮挡。", False, INK)], 13.5, INK, False, BODY)],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "整个项目分成 A、B、C 三个模块。A 模块从 STL 和姿态出发计算 OCS，支持多观测几何；"
         "B 模块用同一套几何和姿态生成光度图像；C 模块再把 OCS、图像以及两者融合用于姿态反演。"
         "这样做的最大好处是，OCS 和图像不是两个割裂的数据源，而是来自同一个物理仿真框架，"
         "所以它们的对比和融合在物理上是一致、可控的。")

# ============================ 第 4 页：复现与验证进展 ============================
s = add_slide()
header(s, "REPRODUCTION · 复现与验证", "前向链路已可复现、可闭合验证", 4)
# 左：三端闭合数据卡 + 要点
bullets(s, Inches(0.6), Inches(1.85), Inches(6.5), Inches(2.5),
    [[("Step 0–8 全流程已完整复现", True)],
     "模块 A / B / C 全部跑通，环境与数据流可追溯",
     [("先验证前向链路，再看反演结果", True)],
     "BRDF、面积、单位、姿态旋转、遮挡逐项闭合验证",
     [("真实卫星 A/B 差异已定性冻结", True)],
     "根因为面元中心采样 vs 像素级可见性语义，属物理差异非代码 bug"],
    size=14.5, gap=8)
# 右上：三端闭合数字
def statcard(x, y, w, h, big, small, col):
    rect(s, x, y, w, h, fill=WHITE, line=col, line_w=1.4, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, y + Inches(0.12), w, Inches(0.62), [(big, 30, col, True, HEAD)], align=PP_ALIGN.CENTER)
    txt(s, x, y + Inches(0.78), w, Inches(0.5), [(small, 12, GRAY, False, BODY)], align=PP_ALIGN.CENTER)
statcard(Inches(7.3), Inches(1.9), Inches(2.6), Inches(1.35), "0.253%", "单平板三端闭合\n平均相对误差", GREEN)
statcard(Inches(10.15), Inches(1.9), Inches(2.6), Inches(1.35), "≤0.25%", "立方体闭合误差\n（凸体自遮挡 0%）", TEAL)
statcard(Inches(7.3), Inches(3.45), Inches(2.6), Inches(1.35), "2.0%", "Full vs Fast\nOCS mean 差异", LGREEN)
statcard(Inches(10.15), Inches(3.45), Inches(2.6), Inches(1.35), "256", "渲染分辨率\n（128 会降 CNN 精度）", GREEN)
# 底部：闭合链
rect(s, Inches(0.6), Inches(5.7), Inches(12.15), Inches(1.05), fill=PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.8), Inches(5.82), Inches(11.8), Inches(0.4),
    [("验证逻辑链（由简到繁，逐层隔离误差来源）", 13, GREEN, True, BODY)])
txt(s, Inches(0.8), Inches(6.22), Inches(11.8), Inches(0.5),
    [([("单平板（无遮挡）", True, INK), ("  →  ", False, TEAL), ("立方体（凸体无自遮挡）", True, INK), ("  →  ", False, TEAL),
       ("L 型（凹体有遮挡）", True, INK), ("  →  ", False, TEAL), ("真实卫星（确认为采样语义差异）", True, GREEN)], 13.5, INK, False, BODY)])
notes(s, "复现这部分已经完整跑通。比较关键的是，我没有只看最终反演结果，而是先把 BRDF、面积、单位、"
         "姿态旋转和遮挡这些前向链路做了闭合验证。单平板和立方体的三端误差都在 0.5% 以内，"
         "所以现在可以认为前向模型没有代码级错误。真实卫星模型里 A、B 两端的差异，"
         "主要来自面元中心采样和像素级可见性语义不同，这是物理离散化差异，已经作为已知现象冻结处理。"
         "另外一个实践经验是，渲染分辨率必须用 256，用 128 会让 CNN 精度明显下降。")

# ============================ 第 5 页：数据与实验设置 ============================
s = add_slide()
header(s, "SETUP · 数据与实验设置", "实验设置与评估口径", 5)
# 左列两块
rect(s, Inches(0.6), Inches(1.85), Inches(6.05), Inches(2.25), fill=CREAM, line=LGREEN, line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.85), Inches(1.98), Inches(5.6), Inches(0.4), [("姿态与数据", 15, GREEN, True, HEAD)])
bullets(s, Inches(0.85), Inches(2.45), Inches(5.6), Inches(1.5),
    [[("yaw–pitch 两轴，roll 固定", True)],
     [("OCS：", True, INK), ("5° 网格 2701 姿态 × 5 观测几何", False, INK)],
     [("图像：", True, INK), ("phase63 单几何，GGX exact BRDF，256", False, INK)]],
    size=13.5, gap=7)
rect(s, Inches(0.6), Inches(4.25), Inches(6.05), Inches(2.45), fill=CREAM, line=TEAL, line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.85), Inches(4.38), Inches(5.6), Inches(0.4), [("划分与指标", 15, GREEN, True, HEAD)])
bullets(s, Inches(0.85), Inches(4.85), Inches(5.6), Inches(1.7),
    [[("主划分：", True, INK), ("10° 网格训练 → 5° 插值测试", False, INK)],
     "测试点不落在训练网格上，比随机划分更严格",
     [("补充：", True, INK), ("random 80/20 split 交叉检查", False, INK)],
     [("指标：", True, INK), ("mean / Hit@5° / Hit@10° / p90 / worst", False, INK)]],
    size=13.5, gap=7, mcolor=TEAL)
# 右：OCS 热图
rect(s, Inches(6.95), Inches(1.85), Inches(5.8), Inches(4.85), fill=WHITE, line=PALE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
pic_fit(s, P['heat'], Inches(7.1), Inches(2.0), Inches(5.5), Inches(4.2))
txt(s, Inches(6.95), Inches(6.18), Inches(5.8), Inches(0.45),
    [("OCS 随 yaw×pitch 的姿态特征（phase63）", 11.5, GRAY, False, BODY)], align=PP_ALIGN.CENTER)
notes(s, "实验设置上，当前主要做 yaw 和 pitch 两个自由度，roll 固定。OCS 部分有 5 个观测几何，"
         "每个几何 2701 个姿态；图像部分目前主要用 phase63 的渲染图。主实验采用 10 度网格训练、"
         "5 度插值测试，这比随机划分更严格，因为测试点不在训练网格上。后面也补了随机 split，"
         "用来回应数据量和划分方式是否特殊的质疑。右边这张图是 OCS 随姿态变化的热图，"
         "可以直观看到 OCS 本身带有很强的姿态结构。")

# ============================ 第 6 页：主结果一 ============================
s = add_slide()
header(s, "RESULT 1 · 基础表现", "OCS、图像与融合的基础反演表现", 6)
# 左：图01
rect(s, Inches(0.55), Inches(1.8), Inches(7.05), Inches(4.25), fill=WHITE, line=PALE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
pic_fit(s, P['fig01'], Inches(0.7), Inches(1.95), Inches(6.75), Inches(3.55))
txt(s, Inches(0.55), Inches(5.6), Inches(7.05), Inches(0.4),
    [("各方法平均姿态误差对比（10°→5° split）", 11.5, GRAY, False, BODY)], align=PP_ALIGN.CENTER)
# 右：数据表（卡片式）
rows = [
    ("OCS MLP · all_raw 45D", "3.98°", "90.7%", "semi-oracle 上界", GREEN),
    ("OCS MLP · per_part_log 30D", "5.91°", "73.8%", "更实用的 OCS", TEAL),
    ("TinyCNN · phase63 图像", "≈12°", "≈26%", "旧图像基线", GRAY),
    ("Feature fusion · per_part", "4.10°", "87.3%", "早期融合 sweet spot", LGREEN),
]
rx = Inches(7.8); rw = Inches(4.95); ry = Inches(1.9); rh = Inches(1.06); rg = Inches(0.12)
for i, (name, mean, hit, note, col) in enumerate(rows):
    y = ry + i * (rh + rg)
    rect(s, rx, y, rw, rh, fill=WHITE, line=col, line_w=1.4, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, rx, y, Inches(0.11), rh, fill=col)
    txt(s, rx + Inches(0.22), y + Inches(0.08), rw - Inches(0.35), Inches(0.42),
        [(name, 13.5, INK, True, BODY)])
    txt(s, rx + Inches(0.22), y + Inches(0.5), Inches(3.4), Inches(0.5),
        [([("mean ", False, GRAY), (mean, True, col), ("   Hit@5° ", False, GRAY), (hit, True, col)], 13, INK, False, BODY)])
    txt(s, rx + Inches(0.22), y + Inches(0.5), rw - Inches(0.4), Inches(0.45),
        [(note, 11, GRAY, False, BODY, True)], align=PP_ALIGN.RIGHT)
notes(s, "从最早的一组主结果看，多几何 OCS 本身已经是很强的姿态信号，all_raw 45 维可以做到大约 4 度，"
         "但它带有比较强的 semi-oracle 性质，含了遮挡率这种实际很难测的量。更实际的 per_part_log 是 5.91 度。"
         "旧的 TinyCNN 图像基线大约 12 度，而 feature fusion 在 per_part_log 条件下能到 4.10 度，"
         "说明中等强度 OCS 和图像之间确实有互补性。这里我不把它说成永远最优，"
         "而是作为早期融合有效性的证据——因为后面 ResNet 的结果会改变这个叙事。")

print("slides 1-6 done")

# ============================ 第 7 页：主结果二 · ResNet 后的重定位 ============================
s = add_slide()
header(s, "RESULT 2 · 论文重定位", "ResNet 结果带来的叙事调整", 7)
# 左：图（ResNet 鲁棒性）
rect(s, Inches(0.55), Inches(1.78), Inches(7.35), Inches(4.1), fill=WHITE, line=PALE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
pic_fit(s, P['resnet'], Inches(0.7), Inches(1.9), Inches(7.05), Inches(3.85))
# 右：要点卡
rect(s, Inches(8.1), Inches(1.78), Inches(4.65), Inches(4.95), fill=CREAM, line=GREEN, line_w=1.3, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(8.3), Inches(1.9), Inches(4.3), Inches(0.45), [("三个关键事实", 15.5, GREEN, True, HEAD)])
bullets(s, Inches(8.32), Inches(2.45), Inches(4.25), Inches(4.1),
    [[("干净仿真图像 = 性能上界", True)],
     [("ResNet-18 clean ", False, INK), ("1.69°", True, GREEN), ("，Hit@5° 97.6%", False, INK)],
     [("强图像模型极度脆弱", True)],
     [("1% 高斯噪声即崩溃至 ", False, INK), ("85.9°", True, RED), ("（Hit@5° 2.2%）", False, INK)],
     [("OCS 对图像退化不敏感", True)],
     [("任何图像噪声下稳定保持 ", False, INK), ("5.91°", True, TEAL)]],
    size=13.5, gap=8)
notes(s, "补了 ResNet-18 以后，叙事发生了一个重要变化。干净仿真图像下，强 CNN 能做到 1.69 度，"
         "说明不能再讲图像只是辅助，也不能讲 fusion 永远最优。但这个结果必须解释为干净仿真图像下的上界。"
         "只要加 1% 的高斯图像噪声，ResNet 就退化到 85 度左右，几乎失效；相反 OCS 完全不受图像噪声影响，"
         "始终是 5.91 度。所以更合理的主线是：图像给出理想条件下的高精度上界，"
         "OCS 提供低成本、可解释、对退化更鲁棒的光度约束，融合的价值是条件性的。这一页是整个汇报的核心转折。")

# ============================ 第 8 页：补充实验回应审稿风险 ============================
s = add_slide()
header(s, "ABLATION · 补充实验", "六项补充实验回应审稿风险", 8)
cards = [
    ("Phase63 公平消融", "单几何 OCS 21.68° → 加图像 6.79°", "单几何下图像补偿更明显", GREEN),
    ("Random split", "Feature fusion per_part 2.13°", "互补性在随机划分下仍成立", TEAL),
    ("BRDF 参数敏感性", "金属 roughness ±20% → OCS 变 30–42%", "非金属部件影响 <5%", LGREEN),
    ("自遮挡 w/ vs w/o", "跨几何遮挡率 60%–78.5%", "自遮挡非装饰模块", GREEN),
    ("Roll 敏感性", "OCS 平均变化 20.3%，max 26.2%", "固定 roll 是明确边界", TEAL),
    ("OCS 噪声鲁棒性", "20% 噪声：OCS 17.25° / 融合 10.96°", "含噪时融合更有价值", LGREEN),
]
cw = Inches(3.95); ch = Inches(1.5); gx = Inches(0.28); gy = Inches(0.22)
x0 = Inches(0.6); y0 = Inches(1.85)
for i, (title, num, note, col) in enumerate(cards):
    r, c = divmod(i, 3)
    x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
    rect(s, x, y, cw, ch, fill=WHITE, line=col, line_w=1.4, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y, cw, Inches(0.1), fill=col)
    txt(s, x + Inches(0.18), y + Inches(0.16), cw - Inches(0.34), Inches(0.4),
        [(title, 14, col, True, HEAD)])
    txt(s, x + Inches(0.18), y + Inches(0.58), cw - Inches(0.34), Inches(0.5),
        [(num, 13, INK, True, BODY)])
    txt(s, x + Inches(0.18), y + Inches(1.04), cw - Inches(0.34), Inches(0.4),
        [(note, 11.5, GRAY, False, BODY, True)])
notes(s, "补充实验主要回应可能的审稿风险。OCS 用 5 几何、图像只用 phase63 是否不公平？"
         "公平消融显示，单几何下图像补偿反而更明显。随机 split 说明结论不是某种划分造成的。"
         "BRDF 敏感性说明材料参数里金属 roughness 最关键，但都在合理公差内。"
         "遮挡分析说明自遮挡影响很大，不是装饰。Roll 敏感性把固定 roll 讲成明确边界而不是回避。"
         "OCS 噪声实验则支持融合在含噪观测下更有价值。这六个实验基本把主要质疑都覆盖了。")

# ============================ 第 9 页：OCS 噪声 → 条件性互补 ============================
s = add_slide()
header(s, "KEY INSIGHT · 条件性互补", "融合价值随观测退化而递增", 9)
rect(s, Inches(0.55), Inches(1.8), Inches(7.55), Inches(4.45), fill=WHITE, line=PALE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
pic_fit(s, P['noise'], Inches(0.7), Inches(1.95), Inches(7.25), Inches(4.15))
rect(s, Inches(8.3), Inches(1.95), Inches(4.45), Inches(4.2), fill=PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(8.5), Inches(2.1), Inches(4.1), Inches(0.45), [("为什么这很重要", 15.5, GREEN, True, HEAD)])
bullets(s, Inches(8.52), Inches(2.65), Inches(4.05), Inches(3.4),
    [[("OCS 与图像误差近乎零相关", True)],
     "r ≈ 0.003，两模态犯不同类型的错误",
     [("噪声越大，图像补偿越关键", True)],
     "增益从 +1.97° 单调升到 +6.29°",
     [("融合不是永远最优，而是条件性", True)],
     "在某一模态退化时提供互补约束"],
    size=13.5, gap=8)
notes(s, "这一页是把上一页的 OCS 噪声实验单独放大讲。横轴是 OCS 测量噪声，蓝线是纯 OCS，"
         "绿线是 OCS 加图像融合，中间绿色区域是图像带来的增益。可以看到一个很干净的规律："
         "噪声越大，图像补偿的增益越大，从 0% 噪声的接近 2 度，一直升到 20% 噪声时的 6.3 度。"
         "再加上 OCS 和图像的误差相关性几乎是零，说明它们犯的是不同类型的错误。"
         "所以融合的正确说法是条件性互补——在某一个模态退化时，另一个能顶上。")

# ============================ 第 10 页：论文写作进展 ============================
s = add_slide()
header(s, "WRITING · 论文进展", "论文写作当前进展", 10)
# 进度条阶段
stages = ["v0.1\n主稿整合", "01 作者确认", "02 引用核验", "03 图表定稿", "04 全文压缩", "05 模拟审稿", "06 投稿材料"]
sw = Inches(1.62); sh = Inches(0.78); sg = Inches(0.10); sx = Inches(0.6); sy = Inches(1.95)
for i, st in enumerate(stages):
    x = sx + i * (sw + sg)
    rect(s, x, sy, sw, sh, fill=GREEN if i == 0 else LGREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, sy + Inches(0.06), sw, sh - Inches(0.1),
        [(st, 11.5, WHITE, True, BODY)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(stages) - 1:
        rect(s, x + sw - Inches(0.01), sy + Inches(0.28), sg + Inches(0.04), Inches(0.22), fill=TEAL, shape=MSO_SHAPE.CHEVRON)
txt(s, Inches(0.6), Inches(2.95), Inches(12.1), Inches(0.4),
    [([("六阶段后整合双线修订均已完成", True, GREEN), ("　→　当前进入 ", False, INK), ("v0.2 前作者统一确认", True, RED)], 15, INK, False, BODY)])
# 左：现状  右：待确认 Blocking
rect(s, Inches(0.6), Inches(3.6), Inches(5.95), Inches(3.05), fill=CREAM, line=LGREEN, line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.82), Inches(3.72), Inches(5.5), Inches(0.4), [("当前状态", 15, GREEN, True, HEAD)])
bullets(s, Inches(0.82), Inches(4.2), Inches(5.5), Inches(2.3),
    ["主稿 v0.1 已整合完成（GPT 底稿 + Claude 组织）",
     "已不在从零写初稿阶段，转为定向修订",
     "暂不生成 v0.2：等 Blocking 项确认后再出新版本"],
    size=13.5, gap=9)
rect(s, Inches(6.8), Inches(3.6), Inches(5.95), Inches(3.05), fill=CREAM, line=RED, line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(7.02), Inches(3.72), Inches(5.5), Inches(0.4), [("待确认 Blocking 项", 15, RED, True, HEAD)])
bullets(s, Inches(7.02), Inches(4.2), Inches(5.5), Inches(2.3),
    ["最终方法定义与关键数值口径", "引用核验、目标期刊", "数据 / 代码共享、作者与声明事实"],
    size=13.5, gap=9, mcolor=RED)
notes(s, "论文写作方面，现在已经不是从零写初稿。v0.1 主稿整合完成后，又做了作者确认、引用、图表、"
         "语言压缩、模拟审稿和投稿材料六个阶段的双线修订，都已完成。现在进入 v0.2 前的统一确认阶段，"
         "卡点在一些需要作者也就是需要我和老师确认的 Blocking 项：最终方法定义、关键数值口径、引用核验、"
         "目标期刊、数据代码共享和作者声明。我的计划是先把这些确认掉，再生成 v0.2。")

# ============================ 第 11 页：阶段性结论 ============================
s = add_slide()
header(s, "CONCLUSION · 阶段性结论", "阶段性结论", 11)
concl = [
    ("01", "前向仿真链路已可复现、可闭合验证", GREEN),
    ("02", "OCS 与图像共享同一 BRDF/几何/姿态框架——物理一致性是项目基础贡献", TEAL),
    ("03", "OCS 是强且可解释的姿态信号，多观测几何非常关键", LGREEN),
    ("04", "干净图像下强 ResNet 是性能上界，但对噪声高度敏感，不可外推到真实观测", GREEN),
    ("05", "融合不是永远最优，而是在中等 OCS 强度或模态退化时提供条件性互补", TEAL),
]
y = Inches(1.95); rh = Inches(0.92); rg = Inches(0.14)
for i, (no, text, col) in enumerate(concl):
    yy = y + i * (rh + rg)
    rect(s, Inches(0.6), yy, Inches(12.15), rh, fill=WHITE, line=col, line_w=1.3, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, Inches(0.6), yy, Inches(0.92), rh, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(0.6), yy, Inches(0.92), rh, [(no, 26, WHITE, True, HEAD)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.75), yy, Inches(10.85), rh, [(text, 15.5, INK, False, BODY)], anchor=MSO_ANCHOR.MIDDLE)
notes(s, "这一阶段可以得出几个比较稳的结论。第一，前向仿真链路已经跑通并通过闭合验证。"
         "第二，项目的核心贡献不是单个模型，而是 OCS 和图像共享同一物理框架。"
         "第三，多几何 OCS 是非常强的姿态信号。第四，ResNet 在干净图像下很强，"
         "但这个结果不能外推到真实退化观测。第五，融合的价值要讲成条件性互补，"
         "尤其在 OCS 信息适中或某个模态退化时更有意义。")

# ============================ 第 12 页：下一步 + 待确认（深绿收尾）============================
s = add_slide()
rect(s, 0, 0, EMU_W, EMU_H, fill=GREEN)
rect(s, 0, 0, EMU_W, Inches(0.22), fill=LGREEN)
rect(s, 0, EMU_H - Inches(0.22), EMU_W, Inches(0.22), fill=TEAL)
txt(s, Inches(0.85), Inches(0.6), Inches(11.5), Inches(0.8), [("下一步计划与待老师确认", 30, WHITE, True, HEAD)])
rect(s, Inches(0.88), Inches(1.45), Inches(2.4), Pt(2.4), fill=LGREEN)
# 左：下一步
rect(s, Inches(0.85), Inches(1.85), Inches(5.7), Inches(4.85), fill=RGBColor(0x05,0x55,0x35), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(1.1), Inches(2.0), Inches(5.2), Inches(0.45), [("下一步计划", 18, LGREEN, True, HEAD)])
bullets(s, Inches(1.1), Inches(2.6), Inches(5.2), Inches(4.0),
    ["清零 v0.2 前的 Blocking 确认项",
     "固定论文主线：clean image 上界 / OCS 鲁棒性 / 条件性融合",
     "完成图表与 caption 最终压缩",
     "进入 v0.2 主稿与投稿材料版本",
     "视时间补高性价比实验：跨 phase 泛化、blur/downsample 退化、ResNet-fusion 鲁棒性"],
    size=14, gap=11, color=WHITE, mcolor=LGREEN)
# 右：待确认
rect(s, Inches(6.8), Inches(1.85), Inches(5.7), Inches(4.85), fill=RGBColor(0x05,0x55,0x35), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(7.05), Inches(2.0), Inches(5.2), Inches(0.45), [("请老师确认", 18, TEAL, True, HEAD)])
bullets(s, Inches(7.05), Inches(2.6), Inches(5.2), Inches(4.0),
    ["是否认可当前主线调整（不再宣传 fusion 绝对最优）",
     "是否需要继续补跨 phase / 图像退化实验",
     "投稿目标：是否先按稳妥 SCI 二区推进",
     "真实 ISAR / 观测数据：纳入主线，还是作为 future work",
     "v0.2 前哪些 Blocking 项优先确认"],
    size=14, gap=11, color=WHITE, mcolor=TEAL)
notes(s, "最后是下一步和需要老师确认的问题。下一步我准备先把 v0.2 前的确认项处理掉，"
         "尤其是方法定义、关键数值和投稿口径，然后把主线固定为 clean image 上界、OCS 鲁棒性和条件性融合，"
         "再进入 v0.2。如果老师觉得有必要，我可以再补一两个高性价比实验。"
         "想请老师确认几个方向性问题：第一，当前主线是否认可；第二，补充实验是否还要继续扩展；"
         "第三，投稿目标和真实数据的定位，是现在纳入主线还是作为 future work 更稳妥。我的汇报就到这里，谢谢老师。")

print("all 12 slides done")
prs.save(OUT)
print("saved ->", OUT)
