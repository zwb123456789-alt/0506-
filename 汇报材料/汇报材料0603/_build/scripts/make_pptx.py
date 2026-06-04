"""Build 20260603 progress-report pptx (18 pages) per spec md section 7.

Style follows reference PDF: deep-green theme, page-header tag, two-column body,
bottom emphasis line, Microsoft YaHei for CJK + Arial for Latin.

Run from project working dir; outputs to ../20260603_项目进展汇报_v2.pptx.
"""

from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- Paths ----------
HERE = Path(__file__).resolve().parent
BUILD = HERE.parent
FIG_DIR = BUILD / "figures"
OUT_PPTX = BUILD.parent / "20260603_项目进展汇报_v2.pptx"

PROJECT = Path(r"d:\我的文件\研究生学术\光学项目\0506新")
EXISTING_FIGS = PROJECT / "结果"

# Verified images (used by name in slide builders)
IMG_PIPELINE        = FIG_DIR / "fig_pipeline_modules.png"
IMG_METHOD_BAR      = FIG_DIR / "fig_method_overview_bar.png"
IMG_RESNET_FUSION   = FIG_DIR / "fig_resnet_fusion_bar.png"
IMG_ROBUSTNESS      = FIG_DIR / "fig_resnet_robustness.png"
IMG_NOISE_CURVE     = FIG_DIR / "fig_noise_fusion_curve.png"
IMG_PHASE63         = FIG_DIR / "fig_phase63_ablation.png"
IMG_ROLL            = FIG_DIR / "fig_roll_sensitivity.png"
IMG_COMPLEMENTARITY = FIG_DIR / "fig_complementarity.png"
IMG_SAMPLING        = FIG_DIR / "fig_sampling_diagnosis.png"

IMG_BRDF_PLANE = EXISTING_FIGS / "BRDF验证" / "plane_batch_20260519_204323" / "fig_plane_batch_compare.png"
IMG_BRDF_CUBE  = EXISTING_FIGS / "BRDF验证" / "cube_20260520_103846" / "fig_cube_compare.png"
IMG_BRDF_L     = EXISTING_FIGS / "BRDF验证" / "L_plate_20260520_103105" / "fig_L_plate_compare.png"

# ---------- Colors ----------
COL_BG          = RGBColor(0xFF, 0xFF, 0xFF)
COL_DEEP_GREEN  = RGBColor(0x0A, 0x5A, 0x3D)   # primary title
COL_MID_GREEN   = RGBColor(0x2F, 0x8F, 0x5A)   # OCS
COL_TEAL        = RGBColor(0x37, 0xB0, 0x95)   # accent (small label)
COL_LIGHT_GREEN = RGBColor(0xD8, 0xED, 0xDC)
COL_SOFT_BG     = RGBColor(0xF3, 0xF7, 0xF4)
COL_ORANGE      = RGBColor(0xC2, 0x49, 0x2A)   # image / warn
COL_ORANGE_BG   = RGBColor(0xFB, 0xE6, 0xDC)
COL_BLUE        = RGBColor(0x1F, 0x6F, 0x9C)   # fusion
COL_BLUE_BG     = RGBColor(0xDA, 0xE7, 0xF0)
COL_TEXT        = RGBColor(0x2B, 0x2B, 0x2B)
COL_TEXT_SOFT   = RGBColor(0x55, 0x55, 0x55)
COL_GREY        = RGBColor(0x7B, 0x7B, 0x7B)
COL_LINE        = RGBColor(0xCC, 0xCC, 0xCC)
COL_RED         = RGBColor(0xC2, 0x32, 0x2D)

# Cover slide deep theme
COL_COVER_BG    = RGBColor(0x0E, 0x4A, 0x36)
COL_COVER_TITLE = RGBColor(0xFF, 0xFF, 0xFF)
COL_COVER_SUB   = RGBColor(0xC9, 0xE8, 0xD3)
COL_COVER_BAR   = RGBColor(0x37, 0xB0, 0x95)

# ---------- Sizes (16:9) ----------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.45)
BODY_TOP = Inches(1.45)
BODY_BOTTOM = Inches(6.85)


# ---------- Helpers ----------
def set_solid_fill(shape, color, line_color=None, line_w=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        if line_w is not None:
            shape.line.width = line_w


def add_rect(slide, left, top, width, height, fill=None, line=None, line_w=None,
             shape_type=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is not None:
        set_solid_fill(shp, fill, line_color=line, line_w=line_w)
    elif line is not None:
        shp.fill.background()
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = line_w
    else:
        shp.fill.background()
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_round_rect(slide, left, top, width, height, fill=None, line=None,
                    line_w=None, corner=0.04):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    # adjust corner
    try:
        shp.adjustments[0] = corner
    except Exception:
        pass
    if fill is not None:
        set_solid_fill(shp, fill, line_color=line, line_w=line_w)
    elif line is not None:
        shp.fill.background()
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = line_w
    else:
        shp.fill.background()
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, runs, *,
             align=PP_ALIGN.LEFT, vanchor=MSO_ANCHOR.TOP,
             auto_size=False, line_spacing=1.15):
    """runs: list of (text, dict(font=, size=, bold=, color=, italic=))
    or list of paragraphs, each paragraph itself a list of runs.
    A 'paragraph' is a list; a 'run' inside is a tuple."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = vanchor
    if not auto_size:
        tf.auto_size = None

    if not runs:
        return tb
    if isinstance(runs[0], tuple):  # single paragraph
        paragraphs = [runs]
    else:
        paragraphs = runs

    for pi, para in enumerate(paragraphs):
        if pi == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for j, item in enumerate(para):
            if isinstance(item, tuple):
                text, attrs = item
            else:
                text, attrs = item, {}
            if j == 0 and p.runs:
                r = p.runs[0]
                r.text = text
            else:
                r = p.add_run()
                r.text = text
            f = r.font
            f.name = attrs.get("font", "Microsoft YaHei")
            if "size" in attrs:
                f.size = Pt(attrs["size"])
            f.bold = bool(attrs.get("bold", False))
            f.italic = bool(attrs.get("italic", False))
            if "color" in attrs:
                f.color.rgb = attrs["color"]
            else:
                f.color.rgb = COL_TEXT
    return tb


def add_picture_fit(slide, image_path, left, top, max_width, max_height,
                    align="center"):
    """Insert picture preserving aspect ratio, fitting within max box."""
    from PIL import Image
    img = Image.open(image_path)
    iw, ih = img.size
    box_w_in = max_width / 914400.0  # EMU to inches
    box_h_in = max_height / 914400.0
    aspect_img = iw / ih
    aspect_box = box_w_in / box_h_in
    if aspect_img >= aspect_box:
        w_in = box_w_in
        h_in = box_w_in / aspect_img
    else:
        h_in = box_h_in
        w_in = box_h_in * aspect_img
    w = Inches(w_in)
    h = Inches(h_in)
    if align == "center":
        l = left + Emu(int((max_width - w) / 2))
        t = top + Emu(int((max_height - h) / 2))
    else:
        l = left
        t = top
    pic = slide.shapes.add_picture(str(image_path), l, t, width=w, height=h)
    return pic


# ---------- Page header / footer ----------
def draw_page_chrome(slide, page_no, header_en, header_cn, title):
    # Top-left small accent bar
    add_rect(slide, MARGIN, Inches(0.42), Inches(0.32), Inches(0.07),
             fill=COL_MID_GREEN)
    # Header tag
    add_text(slide, Inches(0.85), Inches(0.30), Inches(8.0), Inches(0.40),
             [(f"{header_en}  ·  ", {"size": 11, "bold": True, "color": COL_MID_GREEN}),
              (header_cn, {"size": 11, "color": COL_TEXT_SOFT})],
             vanchor=MSO_ANCHOR.MIDDLE)
    # Page number top-right
    add_text(slide, Inches(12.4), Inches(0.30), Inches(0.7), Inches(0.40),
             [(f"{page_no:02d}", {"size": 12, "color": COL_GREY})],
             align=PP_ALIGN.RIGHT, vanchor=MSO_ANCHOR.MIDDLE)
    # Title
    add_text(slide, MARGIN, Inches(0.78), Inches(12.5), Inches(0.65),
             [(title, {"size": 26, "bold": True, "color": COL_DEEP_GREEN})],
             vanchor=MSO_ANCHOR.MIDDLE)
    # Divider line
    add_rect(slide, MARGIN, Inches(1.42), Inches(12.4), Emu(9525),
             fill=COL_LINE)


def add_subtitle(slide, text, top=Inches(1.5)):
    add_text(slide, MARGIN + Inches(0.05), top, Inches(12.4), Inches(0.35),
             [(text, {"size": 13, "color": COL_TEAL, "italic": True})],
             vanchor=MSO_ANCHOR.MIDDLE)


def add_bottom_banner(slide, text, color=COL_DEEP_GREEN, bg=COL_SOFT_BG, edge=COL_LIGHT_GREEN):
    add_round_rect(slide, MARGIN, Inches(6.86), Inches(12.4), Inches(0.45),
                    fill=bg, line=edge, line_w=Pt(0.75), corner=0.4)
    add_text(slide, MARGIN, Inches(6.86), Inches(12.4), Inches(0.45),
             [(text, {"size": 13, "bold": True, "color": color})],
             align=PP_ALIGN.CENTER, vanchor=MSO_ANCHOR.MIDDLE)


def add_bullet_list(slide, left, top, width, height, items, *,
                    bullet_color=COL_MID_GREEN, font_size=12,
                    line_spacing=1.35, head_color=None):
    """items: list of strings; for headed items prefix '##' becomes bold-header."""
    paras = []
    for it in items:
        if it.startswith("## "):
            paras.append([
                ("▍ ", {"size": font_size, "bold": True, "color": bullet_color}),
                (it[3:], {"size": font_size, "bold": True,
                          "color": head_color or COL_DEEP_GREEN}),
            ])
        else:
            paras.append([
                ("▍ ", {"size": font_size, "color": bullet_color}),
                (it, {"size": font_size, "color": COL_TEXT}),
            ])
    add_text(slide, left, top, width, height, paras,
             vanchor=MSO_ANCHOR.TOP, line_spacing=line_spacing)


def add_card(slide, left, top, width, height, title, body_lines, *,
             title_color=COL_DEEP_GREEN, bg=COL_SOFT_BG, edge=COL_LIGHT_GREEN,
             title_size=14, body_size=11.5, body_color=COL_TEXT):
    add_round_rect(slide, left, top, width, height,
                    fill=bg, line=edge, line_w=Pt(0.9), corner=0.06)
    add_text(slide, left + Inches(0.18), top + Inches(0.12),
             width - Inches(0.36), Inches(0.40),
             [(title, {"size": title_size, "bold": True, "color": title_color})])
    paras = []
    for ln in body_lines:
        if ln == "":
            paras.append([("", {"size": body_size})])
            continue
        if ln.startswith("## "):
            paras.append([(ln[3:], {"size": body_size, "bold": True,
                                     "color": title_color})])
        else:
            paras.append([(ln, {"size": body_size, "color": body_color})])
    add_text(slide, left + Inches(0.18), top + Inches(0.55),
             width - Inches(0.36), height - Inches(0.65), paras,
             line_spacing=1.30)


# ---------- Slides ----------
def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COL_COVER_BG)
    # accent bar top-right
    add_rect(slide, Inches(10.8), Inches(0.9), Inches(2.0), Inches(0.12),
             fill=COL_COVER_BAR)
    # main title
    add_text(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.7),
             [[("OCS – 光度图像联合仿真与", {"size": 40, "bold": True, "color": COL_COVER_TITLE})],
              [("空间目标姿态反演  项目进展汇报", {"size": 40, "bold": True, "color": COL_COVER_TITLE})]],
             line_spacing=1.10)
    # subtitle
    add_text(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.5),
             [("统一 BRDF 仿真框架  ·  多模态反演实验结果  ·  下一步论文推进",
               {"size": 17, "color": COL_COVER_SUB})])
    # green divider
    add_rect(slide, Inches(0.8), Inches(5.4), Inches(3.0), Emu(28575),
             fill=COL_COVER_BAR)
    # presenter
    add_text(slide, Inches(0.8), Inches(5.6), Inches(8.0), Inches(0.5),
             [("汇报人:  张文博", {"size": 15, "color": COL_COVER_SUB})])
    add_text(slide, Inches(0.8), Inches(6.05), Inches(8.0), Inches(0.5),
             [("日期:  2026-06-03", {"size": 15, "color": COL_COVER_SUB})])


def slide_2_significance(prs):
    """7.2 现实意义页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COL_BG)
    draw_page_chrome(slide, 2, "SIGNIFICANCE", "现实意义",
                     "现实意义：面向非合作空间目标的光学姿态感知")
    add_subtitle(slide, "观测几何已知时，利用光度响应和图像反推目标姿态",
                 top=Inches(1.50))

    # Core caption banner
    add_round_rect(slide, MARGIN, Inches(1.95), Inches(12.4), Inches(0.78),
                    fill=COL_SOFT_BG, line=COL_LIGHT_GREEN, line_w=Pt(0.8), corner=0.08)
    add_text(slide, MARGIN + Inches(0.25), Inches(1.97), Inches(11.9), Inches(0.78),
             [[("实际地基或空间光学观测中，目标-太阳-探测器几何通常可由星历、观测站位置、观测时间和太阳位置计算得到。",
                {"size": 13, "color": COL_DEEP_GREEN, "bold": True})],
              [("在该几何关系已知的条件下，目标的反射亮度和光度图像会随姿态变化而变化，因此可作为姿态反演依据。",
                {"size": 13, "color": COL_TEXT})]],
             line_spacing=1.30, vanchor=MSO_ANCHOR.MIDDLE)

    # 4 cards 2x2
    cards = [
        ("非合作目标姿态估计",
         "失效卫星、碎片、未知目标或不可通信目标通常缺少姿态遥测，需要依赖外部观测估计姿态。",
         COL_DEEP_GREEN, COL_SOFT_BG, COL_LIGHT_GREEN),
        ("空间态势感知",
         "姿态变化可反映目标翻滚、失稳、异常机动、太阳能板指向变化或结构状态变化。",
         COL_TEAL, RGBColor(0xE5, 0xF3, 0xEE), COL_LIGHT_GREEN),
        ("光学观测可获得",
         "OCS / 光变和光度图像可由地基或空间光学系统获取，是非合作目标长期监测的重要信息来源。",
         COL_MID_GREEN, COL_SOFT_BG, COL_LIGHT_GREEN),
        ("连续状态估计",
         "单次姿态反演结果可作为连续跟踪系统中的观测信息，用于修正和更新目标姿态状态。",
         COL_BLUE, COL_BLUE_BG, RGBColor(0xB5, 0xCE, 0xDC)),
    ]
    w = Inches(6.05)
    h = Inches(1.65)
    gap_x = Inches(0.30)
    gap_y = Inches(0.20)
    x0 = MARGIN
    y0 = Inches(2.92)
    for i, (title, body, tcolor, bg, edge) in enumerate(cards):
        row, col = i // 2, i % 2
        l = x0 + (w + gap_x) * col
        t = y0 + (h + gap_y) * row
        add_round_rect(slide, l, t, w, h, fill=bg, line=edge, line_w=Pt(0.9), corner=0.05)
        add_text(slide, l + Inches(0.20), t + Inches(0.10), w - Inches(0.40), Inches(0.40),
                 [(title, {"size": 15, "bold": True, "color": tcolor})])
        add_text(slide, l + Inches(0.20), t + Inches(0.55), w - Inches(0.40), h - Inches(0.65),
                 [(body, {"size": 12, "color": COL_TEXT})], line_spacing=1.35)

    add_bottom_banner(slide,
                      "目标：将 OCS 光度响应和光度图像转化为空间目标姿态估计中的有效观测约束。")


def slide_3_reliability(prs):
    """7.3 可靠性判断价值页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COL_BG)
    draw_page_chrome(slide, 3, "VALUE", "可靠性判断价值",
                     "实际价值：判断不同观测条件下该信什么结果")
    add_subtitle(slide, "真实应用中，关键不是只输出姿态角，而是判断结果是否可信",
                 top=Inches(1.50))

    # left: 3-line core narrative
    add_round_rect(slide, MARGIN, Inches(1.95), Inches(6.0), Inches(4.85),
                    fill=COL_SOFT_BG, line=COL_LIGHT_GREEN, line_w=Pt(0.8), corner=0.04)
    add_text(slide, MARGIN + Inches(0.25), Inches(2.05),
             Inches(5.5), Inches(0.6),
             [("现实情境", {"size": 16, "bold": True, "color": COL_DEEP_GREEN})])
    add_text(slide, MARGIN + Inches(0.25), Inches(2.55),
             Inches(5.5), Inches(4.0),
             [[("在真实空间目标光学监测中，", {"size": 13, "color": COL_TEXT})],
              [("观测资源和观测质量并不稳定", {"size": 13, "color": COL_DEEP_GREEN, "bold": True})],
              [("", {"size": 6})],
              [("▍ 有时只能获得低维光度 / OCS",
                {"size": 13, "color": COL_TEXT})],
              [("▍ 有时可以获得光度图像，但图像可能模糊、低分辨率或含噪声",
                {"size": 13, "color": COL_TEXT})],
              [("▍ 有时两类观测同时存在，但可靠性并不相同",
                {"size": 13, "color": COL_TEXT})]],
             line_spacing=1.45)

    # right: 4 cards stack
    cards = [
        ("指导观测资源使用", "OCS 在退化图像条件下仍稳定 → 无法获得清晰图像时仍可提供姿态约束。"),
        ("指导融合策略设计", "图像干净时更多依赖图像；图像退化时提高 OCS 权重；OCS 简并时用图像补充。"),
        ("指导结果可信度判断", "不只看模型输出误差，还要判断结果是在理想条件下成立，还是在退化观测下仍可靠。"),
        ("支持后续连续跟踪系统", "OCS 与图像作为姿态滤波器的不同观测来源，可决定每一时刻的更新策略。"),
    ]
    x0 = MARGIN + Inches(6.4)
    w = Inches(6.0)
    h = Inches(1.10)
    gap = Inches(0.10)
    y0 = Inches(1.95)
    colors_cycle = [
        (COL_DEEP_GREEN, COL_SOFT_BG, COL_LIGHT_GREEN),
        (COL_TEAL, RGBColor(0xE5, 0xF3, 0xEE), COL_LIGHT_GREEN),
        (COL_MID_GREEN, COL_SOFT_BG, COL_LIGHT_GREEN),
        (COL_BLUE, COL_BLUE_BG, RGBColor(0xB5, 0xCE, 0xDC)),
    ]
    for i, (t, b) in enumerate(cards):
        tcolor, bg, edge = colors_cycle[i]
        t_top = y0 + (h + gap) * i
        add_round_rect(slide, x0, t_top, w, h, fill=bg, line=edge, line_w=Pt(0.9), corner=0.05)
        add_text(slide, x0 + Inches(0.18), t_top + Inches(0.10),
                 w - Inches(0.36), Inches(0.32),
                 [(t, {"size": 13.5, "bold": True, "color": tcolor})])
        add_text(slide, x0 + Inches(0.18), t_top + Inches(0.45),
                 w - Inches(0.36), h - Inches(0.55),
                 [(b, {"size": 11.5, "color": COL_TEXT})], line_spacing=1.30)

    add_bottom_banner(slide,
                      "本项目的现实价值：为实际光学姿态监测提供\"什么时候用 OCS、什么时候用图像、什么时候融合\"的判断依据。")


def slide_4_complementarity(prs):
    """7.4 OCS 与光度图像互补关系页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COL_BG)
    draw_page_chrome(slide, 4, "COMPLEMENT", "OCS 与图像",
                     "OCS 与光度图像的互补关系")
    add_subtitle(slide, "同源于光学散射，但信息形态不同", top=Inches(1.50))

    # Two big cards
    left_w = Inches(6.0)
    right_w = Inches(6.0)
    top = Inches(1.98)
    height = Inches(4.80)

    # OCS card
    add_round_rect(slide, MARGIN, top, left_w, height,
                    fill=RGBColor(0xE5, 0xF3, 0xEE), line=COL_MID_GREEN, line_w=Pt(1.0),
                    corner=0.04)
    add_text(slide, MARGIN + Inches(0.25), top + Inches(0.12), left_w - Inches(0.5), Inches(0.4),
             [("OCS · 低维光度约束", {"size": 16, "bold": True, "color": COL_DEEP_GREEN})])
    add_text(slide, MARGIN + Inches(0.25), top + Inches(0.60), left_w - Inches(0.5), Inches(4.1),
             [[("优势", {"size": 14, "bold": True, "color": COL_MID_GREEN})],
              [("▍ 可长期、多几何积累", {"size": 12, "color": COL_TEXT})],
              [("▍ 物理含义明确", {"size": 12, "color": COL_TEXT})],
              [("▍ 对图像模糊、分辨率下降和局部成像退化不敏感", {"size": 12, "color": COL_TEXT})],
              [("▍ 多观测几何下可形成稳定姿态约束", {"size": 12, "color": COL_TEXT})],
              [("", {"size": 6})],
              [("局限", {"size": 14, "bold": True, "color": COL_ORANGE})],
              [("▍ 单几何信息量有限", {"size": 12, "color": COL_TEXT})],
              [("▍ 可能存在姿态简并", {"size": 12, "color": COL_TEXT})],
              [("▍ 缺少目标形状和空间结构信息", {"size": 12, "color": COL_TEXT})]],
             line_spacing=1.35)

    # Image card
    x_right = MARGIN + left_w + Inches(0.40)
    add_round_rect(slide, x_right, top, right_w, height,
                    fill=COL_ORANGE_BG, line=COL_ORANGE, line_w=Pt(1.0),
                    corner=0.04)
    add_text(slide, x_right + Inches(0.25), top + Inches(0.12), right_w - Inches(0.5), Inches(0.4),
             [("光度图像 · 空间结构约束", {"size": 16, "bold": True, "color": COL_ORANGE})])
    add_text(slide, x_right + Inches(0.25), top + Inches(0.60), right_w - Inches(0.5), Inches(4.1),
             [[("优势", {"size": 14, "bold": True, "color": COL_MID_GREEN})],
              [("▍ 包含目标投影形状、阴影、亮度分布和部件结构", {"size": 12, "color": COL_TEXT})],
              [("▍ 理想干净图像下姿态信息很强", {"size": 12, "color": COL_TEXT})],
              [("▍ 可补充 OCS 难以区分的对称性姿态", {"size": 12, "color": COL_TEXT})],
              [("", {"size": 6})],
              [("局限", {"size": 14, "bold": True, "color": COL_ORANGE})],
              [("▍ 真实观测易受噪声、湍流、PSF、跟踪误差和低分辨率影响",
                {"size": 12, "color": COL_TEXT})],
              [("▍ clean 图像结果只能代表理想条件上界", {"size": 12, "color": COL_TEXT})],
              [("▍ 合成图像到真实观测存在泛化风险", {"size": 12, "color": COL_TEXT})]],
             line_spacing=1.35)

    add_bottom_banner(slide,
                      "OCS 解决\"低成本、稳定约束\"，图像解决\"空间结构辨识\"，融合解决\"条件变化下的互补判断\"。")


def slide_5_value(prs):
    """7.5 研究价值与关键创新页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COL_BG)
    draw_page_chrome(slide, 5, "CONTRIBUTION", "研究价值与创新",
                     "研究价值：从\"能反演\"到\"知道何时可信\"")
    add_subtitle(slide, "建立 OCS-光度图像联合姿态反演的可验证评估框架",
                 top=Inches(1.50))

    # Top half: 前沿现状/核心问题
    top1 = Inches(1.95)
    add_round_rect(slide, MARGIN, top1, Inches(7.4), Inches(2.05),
                    fill=COL_SOFT_BG, line=COL_LIGHT_GREEN, line_w=Pt(0.9), corner=0.04)
    add_text(slide, MARGIN + Inches(0.20), top1 + Inches(0.10),
             Inches(7.0), Inches(0.4),
             [("前沿现状", {"size": 14, "bold": True, "color": COL_DEEP_GREEN})])
    add_text(slide, MARGIN + Inches(0.20), top1 + Inches(0.50),
             Inches(7.0), Inches(1.6),
             [[("已有多源融合研究：光变 + 测角、光变 + 动力学滤波、",
                {"size": 11.5, "color": COL_TEXT})],
              [("图像 + 惯性/星敏、图像 + 深度/点云/LiDAR 等。",
                {"size": 11.5, "color": COL_TEXT})],
              [("", {"size": 4})],
              [("这些工作说明多源融合有研究基础，但通常不在", {"size": 11.5, "color": COL_TEXT})],
              [("同一套光学散射物理模型下同时生成 OCS 与图像，",
                {"size": 11.5, "bold": True, "color": COL_DEEP_GREEN})],
              [("再比较两者对姿态反演的贡献。", {"size": 11.5, "color": COL_TEXT})]],
             line_spacing=1.35)

    add_round_rect(slide, MARGIN + Inches(7.7), top1, Inches(5.1), Inches(2.05),
                    fill=COL_ORANGE_BG, line=COL_ORANGE, line_w=Pt(0.9), corner=0.04)
    add_text(slide, MARGIN + Inches(7.9), top1 + Inches(0.10),
             Inches(4.7), Inches(0.4),
             [("核心问题", {"size": 14, "bold": True, "color": COL_ORANGE})])
    add_text(slide, MARGIN + Inches(7.9), top1 + Inches(0.55),
             Inches(4.7), Inches(1.5),
             [[("本项目不是只追求一个最低误差结果，",
                {"size": 12, "color": COL_TEXT})],
              [("而是回答一个更接近实际应用的问题：",
                {"size": 12, "color": COL_TEXT})],
              [("", {"size": 4})],
              [("在不同观测质量下，OCS、图像和融合",
                {"size": 12.5, "bold": True, "color": COL_DEEP_GREEN})],
              [("分别什么时候可靠。",
                {"size": 12.5, "bold": True, "color": COL_DEEP_GREEN})]],
             line_spacing=1.35)

    # Bottom: 3 contribution cards
    cards = [
        ("可信前向模型",
         "将 OCS 和光度图像放入同一套 STL 几何、材料、GGX BRDF 和自遮挡模型中，保证两类观测来自同一物理目标。",
         "意义：让后续 OCS 与图像对比具有公平基础。",
         COL_DEEP_GREEN, COL_SOFT_BG, COL_LIGHT_GREEN),
        ("可验证闭环",
         "通过三端闭合、EXR 几何缓冲和采样差异分析，验证仿真链路具有可追溯的物理一致性。",
         "意义：提高仿真结果作为姿态反演基准的可信度。",
         COL_TEAL, RGBColor(0xE5, 0xF3, 0xEE), COL_LIGHT_GREEN),
        ("条件性互补结论",
         "证明融合不是永远最优：理想图像下图像给上界；退化时 OCS 更稳定；一模态失效时另一模态补偿。",
         "意义：为真实观测中如何选择策略提供依据。",
         COL_BLUE, COL_BLUE_BG, RGBColor(0xB5, 0xCE, 0xDC)),
    ]
    top2 = Inches(4.18)
    w = Inches(4.15)
    h = Inches(2.55)
    gap_x = Inches(0.13)
    for i, (t, body, mean, tcolor, bg, edge) in enumerate(cards):
        l = MARGIN + (w + gap_x) * i
        add_round_rect(slide, l, top2, w, h, fill=bg, line=edge, line_w=Pt(0.9), corner=0.05)
        add_text(slide, l + Inches(0.20), top2 + Inches(0.10), w - Inches(0.40), Inches(0.40),
                 [(t, {"size": 14, "bold": True, "color": tcolor})])
        add_text(slide, l + Inches(0.20), top2 + Inches(0.55), w - Inches(0.40), Inches(1.50),
                 [(body, {"size": 11, "color": COL_TEXT})], line_spacing=1.30)
        add_text(slide, l + Inches(0.20), top2 + Inches(2.00), w - Inches(0.40), Inches(0.50),
                 [(mean, {"size": 11, "bold": True, "color": tcolor, "italic": True})],
                 line_spacing=1.25)

    add_bottom_banner(slide, "最终目标是让姿态反演结果不仅\"算得出来\"，而且\"知道什么时候可信\"。")


def slide_6_physical_consistency(prs):
    """7.6 物理一致前向模型必要性页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COL_BG)
    draw_page_chrome(slide, 6, "CONSISTENCY", "物理一致前向模型",
                     "物理一致前向模型的必要性")
    add_subtitle(slide, "为了公平比较 OCS 与图像，而不是比较两套仿真误差",
                 top=Inches(1.50))

    # Left: problem
    top = Inches(1.98)
    add_round_rect(slide, MARGIN, top, Inches(5.85), Inches(4.80),
                    fill=COL_ORANGE_BG, line=COL_ORANGE, line_w=Pt(0.9), corner=0.04)
    add_text(slide, MARGIN + Inches(0.20), top + Inches(0.10),
             Inches(5.5), Inches(0.4),
             [("不统一模型时的问题", {"size": 14, "bold": True, "color": COL_ORANGE})])
    add_text(slide, MARGIN + Inches(0.20), top + Inches(0.55),
             Inches(5.5), Inches(4.1),
             [[("如果 OCS 与图像分别由两套前向模型生成：",
                {"size": 12, "color": COL_TEXT})],
              [("", {"size": 4})],
              [("▍ OCS 端可能使用自定义 BRDF、面元中心采样、",
                {"size": 11.5, "color": COL_TEXT})],
              [("    解析遮挡和一套材料参数",
                {"size": 11.5, "color": COL_TEXT})],
              [("▍ 图像端可能使用 Blender 默认材质、像素级可见性、",
                {"size": 11.5, "color": COL_TEXT})],
              [("    Cycles 内部反射模型和另一套渲染假设",
                {"size": 11.5, "color": COL_TEXT})],
              [("", {"size": 6})],
              [("这时仍可做姿态反演，但很难判断性能差异来自",
                {"size": 12, "color": COL_TEXT})],
              [("观测信息本身，还是来自仿真假设不同。",
                {"size": 12, "bold": True, "color": COL_ORANGE})]],
             line_spacing=1.35)

    # Right: solution
    x_right = MARGIN + Inches(6.15)
    add_round_rect(slide, x_right, top, Inches(6.25), Inches(4.80),
                    fill=COL_SOFT_BG, line=COL_MID_GREEN, line_w=Pt(0.9), corner=0.04)
    add_text(slide, x_right + Inches(0.20), top + Inches(0.10),
             Inches(5.9), Inches(0.4),
             [("一致性的作用 · 本项目做法", {"size": 14, "bold": True, "color": COL_DEEP_GREEN})])
    add_text(slide, x_right + Inches(0.20), top + Inches(0.55),
             Inches(5.9), Inches(4.1),
             [[("控制变量思想 ——", {"size": 12, "bold": True, "color": COL_DEEP_GREEN})],
              [("保持一致：STL 几何、姿态、太阳/探测器几何、",
                {"size": 11.5, "color": COL_TEXT})],
              [("    材料分区、GGX BRDF、自遮挡模型",
                {"size": 11.5, "color": COL_TEXT})],
              [("只改变：观测表达形式 —— OCS 标量 或 像素图像",
                {"size": 11.5, "color": COL_TEXT})],
              [("", {"size": 6})],
              [("本项目做法：", {"size": 12, "bold": True, "color": COL_DEEP_GREEN})],
              [("▍ 剥离 Blender 默认 BRDF，不让黑箱材质决定亮度",
                {"size": 11.5, "color": COL_TEXT})],
              [("▍ Python 显式计算 GGX BRDF",
                {"size": 11.5, "color": COL_TEXT})],
              [("▍ OCS 端和图像后处理端调用同一个 BRDF 函数",
                {"size": 11.5, "color": COL_TEXT})],
              [("▍ Blender 只负责几何可见性和像素级采样",
                {"size": 11.5, "color": COL_TEXT})]],
             line_spacing=1.35)

    add_bottom_banner(slide,
                      "Blender 负责\"看见哪里\"，Python 负责\"如何反光\"，OCS 模块负责\"大规模姿态扫描\"。")


def slide_7_brdf_route(prs):
    """7.7 Blender 几何采样与自定义 BRDF 技术路线页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COL_BG)
    draw_page_chrome(slide, 7, "BRDF ROUTE", "BRDF 技术路线",
                     "Blender 几何采样与自定义 BRDF 的技术路线")
    add_subtitle(slide, "将 Blender 从\"黑箱渲染器\"改造成\"几何采样器\"",
                 top=Inches(1.50))

    # left: 边界表述（防止过度claim）
    top = Inches(1.98)
    add_round_rect(slide, MARGIN, top, Inches(6.0), Inches(2.30),
                    fill=COL_ORANGE_BG, line=COL_ORANGE, line_w=Pt(0.9), corner=0.04)
    add_text(slide, MARGIN + Inches(0.18), top + Inches(0.10),
             Inches(5.7), Inches(0.40),
             [("重要表述边界", {"size": 13.5, "bold": True, "color": COL_ORANGE})])
    add_text(slide, MARGIN + Inches(0.18), top + Inches(0.50),
             Inches(5.7), Inches(1.85),
             [[("不写\"本文 GGX 比 Blender BRDF 更先进\"。",
                {"size": 11, "color": COL_TEXT})],
              [("Blender Principled BSDF 本身也含 GGX / 多重散射 GGX，",
                {"size": 11, "color": COL_TEXT})],
              [("技术上并不落后，但面向视觉渲染综合材质模型。",
                {"size": 11, "color": COL_TEXT})],
              [("", {"size": 3})],
              [("自定义 GGX 的优势不是\"绝对更先进\"，",
                {"size": 11, "color": COL_TEXT})],
              [("而是\"更适合本文任务\"：公式显式、两端共享、",
                {"size": 11, "bold": True, "color": COL_DEEP_GREEN})],
              [("可控、可验证、可替换、可做敏感性分析。",
                {"size": 11, "bold": True, "color": COL_DEEP_GREEN})]],
             line_spacing=1.30)

    # right: advantages 5 lines
    x_right = MARGIN + Inches(6.30)
    add_round_rect(slide, x_right, top, Inches(6.10), Inches(2.30),
                    fill=COL_SOFT_BG, line=COL_MID_GREEN, line_w=Pt(0.9), corner=0.04)
    add_text(slide, x_right + Inches(0.18), top + Inches(0.10),
             Inches(5.8), Inches(0.40),
             [("自定义 BRDF 的技术优势", {"size": 13.5, "bold": True, "color": COL_DEEP_GREEN})])
    advs = [
        ("公式显式", "GGX/Cook-Torrance 的 D/G/F 可逐项说明，写入方法部分"),
        ("两端共享", "OCS 端与图像后处理端调用同一个 BRDF 函数"),
        ("参数可控", "三材料统一参数库，便于敏感性分析"),
        ("可更新", "可扩展各向异性 / 多光谱 / 材料老化模型"),
        ("计算高效", "同一批 EXR 反复用于不同 BRDF/参数实验"),
    ]
    paras = []
    for k, v in advs:
        paras.append([(f"▍ {k}：", {"size": 11.5, "bold": True, "color": COL_MID_GREEN}),
                      (v, {"size": 11.5, "color": COL_TEXT})])
    add_text(slide, x_right + Inches(0.18), top + Inches(0.50),
             Inches(5.8), Inches(1.85), paras, line_spacing=1.30)

    # bottom: 文献支撑
    top2 = Inches(4.40)
    add_round_rect(slide, MARGIN, top2, Inches(12.4), Inches(2.40),
                    fill=COL_SOFT_BG, line=COL_LIGHT_GREEN, line_w=Pt(0.9), corner=0.04)
    add_text(slide, MARGIN + Inches(0.20), top2 + Inches(0.10),
             Inches(12.0), Inches(0.40),
             [("文献支撑与 OCS 借鉴 Blender 像素级采样的原因",
               {"size": 14, "bold": True, "color": COL_DEEP_GREEN})])
    add_text(slide, MARGIN + Inches(0.20), top2 + Inches(0.55),
             Inches(12.0), Inches(2.0),
             [[("Cook & Torrance, 1982 · Walter et al., 2007 (GGX) · Burley, 2012 (Disney Principled BSDF) · Blender Manual: Principled BSDF",
                {"size": 11, "italic": True, "color": COL_TEXT_SOFT})],
              [("", {"size": 4})],
              [("结论：本文不否定 Blender 内置材质，而是将其从最终亮度计算黑箱中剥离出来，",
                {"size": 12, "color": COL_TEXT})],
              [("改用显式 GGX/Cook-Torrance 作为 OCS 与图像共同的、可审计的散射模型。",
                {"size": 12, "bold": True, "color": COL_DEEP_GREEN})],
              [("", {"size": 4})],
              [("OCS 借鉴 Blender 像素级采样的原因：",
                {"size": 12, "bold": True, "color": COL_TEAL})],
              [("面元中心采样速度快但对 GGX 金属窄镜面峰可能漏掉局部高光；Blender 光栅化提供更细可见区采样。",
                {"size": 12, "color": COL_TEXT})]],
             line_spacing=1.30)

    add_bottom_banner(slide,
                      "直接用 Blender 能生成图像；拆解 Blender 并自定义 BRDF，才能生成可验证、可解释、可用于姿态反演的物理观测数据。")


def slide_8_sampling_diag(prs):
    """7.8 面元中心采样与像素级采样差异诊断页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COL_BG)
    draw_page_chrome(slide, 8, "DIAGNOSIS", "采样差异诊断",
                     "面元中心采样与像素级采样的差异诊断")
    add_subtitle(slide, "真实卫星复杂几何下，面元中心采样会低估窄镜面峰",
                 top=Inches(1.50))

    # Left: figure
    top = Inches(1.95)
    add_round_rect(slide, MARGIN, top, Inches(7.4), Inches(4.85),
                    fill=COL_BG, line=COL_LIGHT_GREEN, line_w=Pt(0.9), corner=0.03)
    add_picture_fit(slide, IMG_SAMPLING,
                    MARGIN + Inches(0.10), top + Inches(0.10),
                    Inches(7.20), Inches(4.65))

    # Right: key findings card
    x_right = MARGIN + Inches(7.7)
    w = Inches(5.10)
    add_round_rect(slide, x_right, top, w, Inches(4.85),
                    fill=COL_SOFT_BG, line=COL_LIGHT_GREEN, line_w=Pt(0.9), corner=0.04)
    add_text(slide, x_right + Inches(0.20), top + Inches(0.10),
             w - Inches(0.40), Inches(0.40),
             [("典型姿态 yaw=150° / pitch=−80°", {"size": 13.5, "bold": True, "color": COL_DEEP_GREEN})])
    add_text(slide, x_right + Inches(0.20), top + Inches(0.55),
             w - Inches(0.40), Inches(1.8),
             [[("▍ A 面元中心 with_occ ≈ 0.0163", {"size": 11.5, "color": COL_TEXT})],
              [("▍ A 面元中心 no_occ ≈ 0.0766", {"size": 11.5, "color": COL_TEXT})],
              [("▍ B 像素级 ≈ 0.1711", {"size": 11.5, "color": COL_TEXT})],
              [("", {"size": 4})],
              [("差异并非 BRDF 公式 / 法线坐标 / 网格精度问题。",
                {"size": 11.5, "italic": True, "color": COL_ORANGE})]],
             line_spacing=1.30)

    add_text(slide, x_right + Inches(0.20), top + Inches(2.40),
             w - Inches(0.40), Inches(2.30),
             [[("关键诊断", {"size": 13, "bold": True, "color": COL_TEAL})],
              [("① 几何精度假说被排除", {"size": 11, "bold": True, "color": COL_TEXT})],
              [("   A_fast vs A_full no_occ 差仅 0.2%", {"size": 10.5, "color": COL_TEXT_SOFT})],
              [("② 窄镜面峰采样问题被确认", {"size": 11, "bold": True, "color": COL_TEXT})],
              [("   A 镜面贡献≈0；B 镜面贡献占 87%", {"size": 10.5, "color": COL_TEXT_SOFT})],
              [("③ 可见性语义存在差异", {"size": 11, "bold": True, "color": COL_TEXT})],
              [("   diffuse-only 下 A/B 仍差 26%", {"size": 10.5, "color": COL_TEXT_SOFT})]],
             line_spacing=1.30)

    add_bottom_banner(slide,
                      "面元中心采样适合大规模扫描；像素级采样可作为 OCS-图像一致性分析的关键补充。")


def slide_9_three_end_closure(prs):
    """7.9 三端闭合验证与采样路线选择页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COL_BG)
    draw_page_chrome(slide, 9, "CLOSURE", "三端闭合验证",
                     "三端闭合验证与采样路线选择")
    add_subtitle(slide, "先排除公式和单位问题，再定位采样差异", top=Inches(1.50))

    # row 1: 3 BRDF validation pics
    top = Inches(1.95)
    pic_h = Inches(2.55)
    w_each = Inches(4.05)
    gap = Inches(0.10)
    pics = [
        (IMG_BRDF_PLANE, "单平板三端闭合", "mean rel_err ≈ 0.253%"),
        (IMG_BRDF_CUBE,  "立方体三端闭合",   "B/an ≤ 0.25% (凸体自遮挡 0%)"),
        (IMG_BRDF_L,     "L 型双平板",       "中等角度 A_with/B ≈ 1.0"),
    ]
    for i, (path, label, sub) in enumerate(pics):
        l = MARGIN + (w_each + gap) * i
        add_round_rect(slide, l, top, w_each, pic_h,
                        fill=COL_BG, line=COL_LIGHT_GREEN, line_w=Pt(0.7), corner=0.03)
        try:
            add_picture_fit(slide, path, l + Inches(0.08), top + Inches(0.08),
                            w_each - Inches(0.16), pic_h - Inches(0.16))
        except Exception:
            add_text(slide, l + Inches(0.5), top + Inches(1.0),
                     w_each - Inches(1.0), Inches(0.6),
                     [("(图缺失)", {"size": 13, "color": COL_GREY})],
                     align=PP_ALIGN.CENTER)

    # labels under each pic
    for i, (path, label, sub) in enumerate(pics):
        l = MARGIN + (w_each + gap) * i
        add_text(slide, l, top + pic_h + Inches(0.05), w_each, Inches(0.32),
                 [(label, {"size": 12, "bold": True, "color": COL_DEEP_GREEN})],
                 align=PP_ALIGN.CENTER)
        add_text(slide, l, top + pic_h + Inches(0.35), w_each, Inches(0.30),
                 [(sub, {"size": 10.5, "color": COL_TEAL, "italic": True})],
                 align=PP_ALIGN.CENTER)

    # row 2: 4 key numbers
    top2 = Inches(5.10)
    metrics = [
        ("0.253%", "单平板\n三端平均相对误差", COL_DEEP_GREEN),
        ("≤0.25%", "立方体\nB/analytical 误差", COL_TEAL),
        ("≈1.0",   "L 型 A_with/B\n中等角度遮挡一致", COL_MID_GREEN),
        ("face-center\nvs pixel", "复杂几何下\n可见性语义差异", COL_ORANGE),
    ]
    w_card = Inches(3.05)
    gap2 = Inches(0.13)
    for i, (val, label, color) in enumerate(metrics):
        l = MARGIN + (w_card + gap2) * i
        add_round_rect(slide, l, top2, w_card, Inches(1.65),
                        fill=COL_SOFT_BG, line=color, line_w=Pt(1.1), corner=0.06)
        add_text(slide, l + Inches(0.1), top2 + Inches(0.10),
                 w_card - Inches(0.2), Inches(0.85),
                 [(val, {"size": 22, "bold": True, "color": color})],
                 align=PP_ALIGN.CENTER, vanchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, l + Inches(0.1), top2 + Inches(0.95),
                 w_card - Inches(0.2), Inches(0.65),
                 [[(label.split("\n")[0], {"size": 11, "color": COL_TEXT_SOFT})],
                  [(label.split("\n")[1], {"size": 11, "color": COL_TEXT_SOFT})]],
                 align=PP_ALIGN.CENTER, line_spacing=1.2)

    add_bottom_banner(slide,
                      "像素级采样不是替代全部 OCS 扫描，而是在复杂高光和可见性场景下提高前向模型的一致性与可信度。")


def slide_10_route_status(prs):
    """7.10 路线结果页1：三模块工作路线与当前完成状态"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COL_BG)
    draw_page_chrome(slide, 10, "PIPELINE", "总体技术流程",
                     "三模块工作路线与当前完成状态")

    # Insert pipeline figure (covers most of body)
    top = Inches(1.65)
    add_picture_fit(slide, IMG_PIPELINE,
                    MARGIN, top, Inches(12.4), Inches(3.6))

    # bottom row: 3 completion-status cards
    top2 = Inches(5.30)
    cards = [
        ("模块 A · OCS 计算",
         ["▍ 5° 网格 2701 姿态 × 5 观测几何", "▍ 总耗时 28.7 min (OPTIX GPU)",
          "▍ run_20260520_162831 已落地"],
         COL_MID_GREEN, COL_SOFT_BG),
        ("模块 B · Blender 渲染",
         ["▍ phase63 2701 帧 EXR + 后处理 PNG", "▍ 0.32 s/帧 OPTIX GPU",
          "▍ run_20260521_phase63_ggx 已落地"],
         COL_TEAL, RGBColor(0xE5, 0xF3, 0xEE)),
        ("模块 C · 姿态反演",
         ["▍ OCS-only / image-only / 融合 全部跑通", "▍ ResNet 主线已收敛",
          "▍ paper_summary run_20260528_162418"],
         COL_DEEP_GREEN, COL_SOFT_BG),
    ]
    w = Inches(4.08)
    gap = Inches(0.10)
    for i, (title, items, color, bg) in enumerate(cards):
        l = MARGIN + (w + gap) * i
        add_round_rect(slide, l, top2, w, Inches(1.55),
                        fill=bg, line=color, line_w=Pt(1.0), corner=0.05)
        add_text(slide, l + Inches(0.15), top2 + Inches(0.08), w - Inches(0.3), Inches(0.36),
                 [(title, {"size": 13.5, "bold": True, "color": color})])
        paras = []
        for it in items:
            paras.append([(it, {"size": 11, "color": COL_TEXT})])
        add_text(slide, l + Inches(0.15), top2 + Inches(0.45), w - Inches(0.3), Inches(1.05),
                 paras, line_spacing=1.30)

    add_bottom_banner(slide, "三模块均已落地；下一步关注 v0.2 之前的图表与口径锁定。")


def slide_11_forward_validation(prs):
    """前向模型验证结果"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COL_BG)
    draw_page_chrome(slide, 11, "VALIDATION", "前向模型验证",
                     "前向链路已可复现、可闭合验证")
    add_subtitle(slide, "单平板 → 立方体 → L 型 → 真实卫星：由简到繁，逐层隔离误差来源",
                 top=Inches(1.50))

    # 4 cards row
    top = Inches(2.00)
    cards = [
        ("Step 0–8 复现",
         "模块 A / B / C 全部跑通；环境、Blender 路径、Python 依赖、数据流均可追溯。",
         COL_DEEP_GREEN),
        ("BRDF 三端闭合",
         "解析 / Python OCS / Blender EXR 后处理三端在单平板/立方体上 ≤ 0.25%。",
         COL_TEAL),
        ("Full vs Fast",
         "fast 精度 OCS mean 与 full 差异约 2.0%，论文期切 full 即可。",
         COL_MID_GREEN),
        ("真实卫星 A/B 冻结",
         "差异根因 = 面元中心 vs 像素级可见性语义，属物理建模差异而非代码 bug。",
         COL_ORANGE),
    ]
    w = Inches(3.05)
    gap = Inches(0.10)
    for i, (t, b, color) in enumerate(cards):
        l = MARGIN + (w + gap) * i
        add_round_rect(slide, l, top, w, Inches(1.85),
                        fill=COL_SOFT_BG, line=color, line_w=Pt(1.0), corner=0.05)
        add_text(slide, l + Inches(0.15), top + Inches(0.10), w - Inches(0.3), Inches(0.36),
                 [(t, {"size": 13, "bold": True, "color": color})])
        add_text(slide, l + Inches(0.15), top + Inches(0.50), w - Inches(0.3), Inches(1.30),
                 [(b, {"size": 11, "color": COL_TEXT})], line_spacing=1.32)

    # 4 key-number tiles (with sampling diagnosis pic on right)
    top2 = Inches(4.05)
    add_picture_fit(slide, IMG_SAMPLING, MARGIN, top2,
                    Inches(12.4), Inches(2.75))

    add_bottom_banner(slide,
                      "前向模型已具备\"先验证、后反演\"的可信基础；后续反演结果具有公平的对比意义。")


def slide_12_methods_definition(prs):
    """反演对比1：反演方法定义"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COL_BG)
    draw_page_chrome(slide, 12, "METHODS", "反演方法定义",
                     "四类反演方法：输入 · 网络 · 输出")

    # 4 method cards 2x2
    top = Inches(1.65)
    w = Inches(6.05)
    h = Inches(2.50)
    gap_x = Inches(0.30)
    gap_y = Inches(0.20)
    methods = [
        ("仅 OCS",
         "输入：多观测几何 OCS 特征 (concat5 30D / 45D)",
         "网络：MLP 128→128→64 SiLU LayerNorm",
         "输出：[sin yaw, cos yaw, sin pitch, cos pitch]",
         "优势：低维、可解释、多几何稳定； 局限：可能姿态简并。",
         COL_MID_GREEN, COL_SOFT_BG),
        ("仅图像",
         "输入：phase63 单几何光度图像 1×128×128 (log1p)",
         "网络：TinyCNN 106k / ResNet-18 11M",
         "输出：与 OCS 端一致的 sin/cos 编码",
         "优势：包含形状、阴影、空间亮度； 局限：对真实图像退化敏感。",
         COL_ORANGE, COL_ORANGE_BG),
        ("预测级融合 (Late)",
         "输入：OCS 与图像模型各自输出的 sin/cos 向量",
         "方法：vec_fused = β·vec_ocs + (1−β)·vec_img",
         "调参：β sweep 0:0.01:1",
         "优点：简单； 局限：不能学习跨模态深层关系。",
         COL_BLUE, COL_BLUE_BG),
        ("特征级融合 (Feature)",
         "输入：OCS 特征 + 图像同时输入",
         "网络：OCS-MLP 64D ⊕ ImageBranch 64D → FusionHead",
         "训练：end-to-end，5 seeds 平均",
         "优点：可学互补信息； 局限：对样本量/训练稳定性更敏感。",
         COL_DEEP_GREEN, COL_SOFT_BG),
    ]
    for i, (t, l1, l2, l3, l4, color, bg) in enumerate(methods):
        row, col = i // 2, i % 2
        x = MARGIN + (w + gap_x) * col
        y = top + (h + gap_y) * row
        add_round_rect(slide, x, y, w, h, fill=bg, line=color, line_w=Pt(1.0), corner=0.04)
        add_text(slide, x + Inches(0.20), y + Inches(0.10), w - Inches(0.4), Inches(0.38),
                 [(t, {"size": 15, "bold": True, "color": color})])
        body = [
            [(l1, {"size": 11.5, "color": COL_TEXT})],
            [(l2, {"size": 11.5, "color": COL_TEXT})],
            [(l3, {"size": 11.5, "color": COL_TEXT})],
            [("", {"size": 4})],
            [(l4, {"size": 11.5, "italic": True, "color": COL_TEXT_SOFT})],
        ]
        add_text(slide, x + Inches(0.20), y + Inches(0.55), w - Inches(0.4), h - Inches(0.6),
                 body, line_spacing=1.30)

    add_bottom_banner(slide,
                      "四类方法共享相同 OCS / 图像数据源和评估口径，差异仅在\"如何利用观测信息\"。")


def slide_13_main_results(prs):
    """反演对比2：主反演结果对比"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COL_BG)
    draw_page_chrome(slide, 13, "RESULT MAIN", "主反演结果",
                     "OCS、图像与融合的主反演表现")
    add_subtitle(slide, "10°→5° 插值划分；测试点不在训练网格上",
                 top=Inches(1.50))

    # large figure
    add_picture_fit(slide, IMG_METHOD_BAR, MARGIN, Inches(1.92),
                    Inches(8.4), Inches(4.85))

    # right cards
    x_right = MARGIN + Inches(8.70)
    w = Inches(3.70)
    items = [
        ("OCS MLP all_raw 45D",
         "mean 3.98°  Hit5 90.7%  semi-oracle 上界",
         COL_MID_GREEN, COL_SOFT_BG),
        ("ResNet-18 仅图像 clean",
         "mean 1.69°  Hit5 97.6%  干净图像上界",
         COL_ORANGE, COL_ORANGE_BG),
        ("Feature fusion ResNet+OCS",
         "mean 1.47°  Hit5 99.7%  worst 6.62°",
         COL_DEEP_GREEN, COL_LIGHT_GREEN),
        ("Late fusion OCS+CNN",
         "mean 5.03°  β=0.96 时受 OCS 主导",
         COL_BLUE, COL_BLUE_BG),
    ]
    top_card = Inches(1.95)
    h = Inches(1.13)
    gap = Inches(0.10)
    for i, (t, b, color, bg) in enumerate(items):
        y = top_card + (h + gap) * i
        add_round_rect(slide, x_right, y, w, h, fill=bg, line=color, line_w=Pt(0.9), corner=0.05)
        add_text(slide, x_right + Inches(0.15), y + Inches(0.08),
                 w - Inches(0.3), Inches(0.36),
                 [(t, {"size": 12, "bold": True, "color": color})])
        add_text(slide, x_right + Inches(0.15), y + Inches(0.46),
                 w - Inches(0.3), h - Inches(0.5),
                 [(b, {"size": 11, "color": COL_TEXT})], line_spacing=1.25)

    add_bottom_banner(slide,
                      "干净图像下 ResNet 即达性能上界；融合 OCS 进一步将 worst-case 由 9.9° 压到 6.6°。")


def slide_14_resnet_repositioning(prs):
    """反演对比3：ResNet 后主线调整"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COL_BG)
    draw_page_chrome(slide, 14, "REPOSITIONING", "主线调整",
                     "ResNet 结果带来的主线调整")
    add_subtitle(slide, "强图像模型 = 干净仿真上界；OCS 仍是真实观测下的鲁棒约束",
                 top=Inches(1.50))

    # left: fusion bar chart
    add_picture_fit(slide, IMG_RESNET_FUSION,
                    MARGIN, Inches(1.95), Inches(7.6), Inches(4.85))

    # right: 3 key findings
    x_right = MARGIN + Inches(7.85)
    w = Inches(4.55)
    items = [
        ("干净仿真图像 = 性能上界",
         "ResNet-18 clean mean 1.69°、Hit@5° 97.6%；显著低于 OCS-only 5.91°。",
         COL_MID_GREEN, COL_SOFT_BG),
        ("OCS 仍带来 worst-case 改善",
         "ResNet+OCS per_part 30D 将 worst 9.9° → 6.6° (−33%)，mean 1.69° → 1.47°。",
         COL_DEEP_GREEN, COL_LIGHT_GREEN),
        ("旧主线需要调整",
         "不再宣传\"融合永远最优\"；重写为 clean image 上界 + OCS 鲁棒约束 + 条件性融合。",
         COL_ORANGE, COL_ORANGE_BG),
    ]
    top = Inches(1.98)
    h = Inches(1.55)
    gap = Inches(0.13)
    for i, (t, b, color, bg) in enumerate(items):
        y = top + (h + gap) * i
        add_round_rect(slide, x_right, y, w, h, fill=bg, line=color, line_w=Pt(1.0), corner=0.05)
        add_text(slide, x_right + Inches(0.18), y + Inches(0.10),
                 w - Inches(0.36), Inches(0.4),
                 [(t, {"size": 13.5, "bold": True, "color": color})])
        add_text(slide, x_right + Inches(0.18), y + Inches(0.50),
                 w - Inches(0.36), h - Inches(0.55),
                 [(b, {"size": 11.5, "color": COL_TEXT})], line_spacing=1.32)

    add_bottom_banner(slide,
                      "ResNet 拉高图像上界 → 论文叙事必须从\"融合最优\"调整为\"条件性互补\"。")


def slide_15_robustness_image(prs):
    """鲁棒性1：图像退化结果"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COL_BG)
    draw_page_chrome(slide, 15, "ROBUSTNESS I", "图像退化鲁棒性",
                     "图像鲁棒性：clean 上界并不外推到真实观测")
    add_subtitle(slide, "1% 高斯噪声即让 ResNet 从 1.69° 崩到 85.85°；OCS-only 不受影响",
                 top=Inches(1.50))

    # large fig
    add_picture_fit(slide, IMG_ROBUSTNESS,
                    MARGIN, Inches(1.95), Inches(8.4), Inches(4.85))

    # right: 3 findings
    x_right = MARGIN + Inches(8.70)
    w = Inches(3.70)
    items = [
        ("clean 仿真 = 上界",
         "ResNet-18 在仿真图像上 1.69°；这是图像反演的能力天花板。",
         COL_MID_GREEN, COL_SOFT_BG),
        ("强 CNN 极度脆弱",
         "1%/3%/5%/10% 高斯噪声均退化到 85°+；Hit@5° 跌至 1–2%。",
         COL_ORANGE, COL_ORANGE_BG),
        ("OCS 对像素退化不敏感",
         "图像噪声/亮度变换不会影响 OCS 数值（同一光度积分）；保持 5.91°。",
         COL_DEEP_GREEN, COL_LIGHT_GREEN),
    ]
    top = Inches(1.98)
    h = Inches(1.55)
    gap = Inches(0.13)
    for i, (t, b, color, bg) in enumerate(items):
        y = top + (h + gap) * i
        add_round_rect(slide, x_right, y, w, h, fill=bg, line=color, line_w=Pt(1.0), corner=0.05)
        add_text(slide, x_right + Inches(0.16), y + Inches(0.10),
                 w - Inches(0.32), Inches(0.4),
                 [(t, {"size": 13, "bold": True, "color": color})])
        add_text(slide, x_right + Inches(0.16), y + Inches(0.50),
                 w - Inches(0.32), h - Inches(0.55),
                 [(b, {"size": 11, "color": COL_TEXT})], line_spacing=1.32)

    add_bottom_banner(slide,
                      "clean image 性能不可直接外推到真实外场观测；OCS 是真实退化条件下的鲁棒后盾。")


def slide_16_robustness_ocs(prs):
    """鲁棒性2：OCS 噪声与融合补偿"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COL_BG)
    draw_page_chrome(slide, 16, "ROBUSTNESS II", "OCS 噪声与融合补偿",
                     "条件性互补：融合价值随观测退化而递增")
    add_subtitle(slide, "OCS 越糟，图像补偿越大；r ≈ 0.003 → 完全互补",
                 top=Inches(1.50))

    # left fig
    add_picture_fit(slide, IMG_NOISE_CURVE,
                    MARGIN, Inches(1.95), Inches(7.8), Inches(4.85))

    # right: complementarity scatter + bullets
    x_right = MARGIN + Inches(7.95)
    w = Inches(4.45)
    add_picture_fit(slide, IMG_COMPLEMENTARITY,
                    x_right, Inches(1.95), w, Inches(2.45))
    # bullets
    add_round_rect(slide, x_right, Inches(4.55), w, Inches(2.25),
                    fill=COL_SOFT_BG, line=COL_LIGHT_GREEN, line_w=Pt(0.9), corner=0.04)
    add_text(slide, x_right + Inches(0.15), Inches(4.65),
             w - Inches(0.30), Inches(0.4),
             [("核心结论", {"size": 13.5, "bold": True, "color": COL_DEEP_GREEN})])
    add_text(slide, x_right + Inches(0.15), Inches(5.05),
             w - Inches(0.30), Inches(1.75),
             [[("▍ OCS 与图像误差近乎零相关 (r≈0.003)",
                {"size": 11, "color": COL_TEXT})],
              [("▍ 噪声越大，图像补偿越关键：+1.97° → +6.29°",
                {"size": 11, "color": COL_TEXT})],
              [("▍ 融合不是永远最优，而是条件性",
                {"size": 11, "bold": True, "color": COL_DEEP_GREEN})],
              [("▍ 64.9% 样本融合优于两种单模态",
                {"size": 11, "color": COL_TEXT})]],
             line_spacing=1.35)

    add_bottom_banner(slide,
                      "融合的价值随观测退化单调递增；这才是\"OCS+图像\"在真实场景的应用空间。")


def slide_17_ablation_six(prs):
    """补充实验页：审稿风险防御"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COL_BG)
    draw_page_chrome(slide, 17, "ABLATION", "补充实验",
                     "六项补充实验回应审稿风险")
    add_subtitle(slide, "phase63 公平消融 / random split / BRDF / 遮挡 / roll / OCS 噪声",
                 top=Inches(1.50))

    # left: phase63 ablation fig + roll fig
    add_picture_fit(slide, IMG_PHASE63,
                    MARGIN, Inches(1.95), Inches(5.9), Inches(2.30))
    add_picture_fit(slide, IMG_ROLL,
                    MARGIN, Inches(4.35), Inches(5.9), Inches(2.45))

    # right: 6 small cards 3x2
    x_right = MARGIN + Inches(6.10)
    cards = [
        ("Phase63 公平消融", "单几何 OCS 21.68° → 加图像 6.79° (−14.89°)",
         COL_MID_GREEN, COL_SOFT_BG),
        ("Random split", "Feature fusion per_part 2.13°，互补性在随机划分仍成立",
         COL_TEAL, RGBColor(0xE5, 0xF3, 0xEE)),
        ("BRDF 参数敏感性", "金属 roughness ±20% → OCS 变 30–42%；非金属 <5%",
         COL_DEEP_GREEN, COL_SOFT_BG),
        ("自遮挡 w/ vs w/o", "跨几何遮挡率 60%–78.5%，自遮挡非装饰模块",
         COL_BLUE, COL_BLUE_BG),
        ("Roll 敏感性", "OCS 平均变化 20.3%，max 26.2%；固定 roll 是论文边界",
         COL_ORANGE, COL_ORANGE_BG),
        ("OCS 噪声鲁棒性", "20% 噪声：OCS 17.25° / 融合 10.96°；含噪时融合更有价值",
         COL_DEEP_GREEN, COL_LIGHT_GREEN),
    ]
    w = Inches(3.13)
    h = Inches(1.55)
    gap_x = Inches(0.08)
    gap_y = Inches(0.10)
    for i, (t, b, color, bg) in enumerate(cards):
        col = i % 2
        row = i // 2
        l = x_right + (w + gap_x) * col
        y = Inches(1.95) + (h + gap_y) * row
        add_round_rect(slide, l, y, w, h, fill=bg, line=color, line_w=Pt(0.9), corner=0.06)
        add_text(slide, l + Inches(0.14), y + Inches(0.10),
                 w - Inches(0.28), Inches(0.36),
                 [(t, {"size": 12, "bold": True, "color": color})])
        add_text(slide, l + Inches(0.14), y + Inches(0.48),
                 w - Inches(0.28), h - Inches(0.55),
                 [(b, {"size": 10.5, "color": COL_TEXT})], line_spacing=1.30)

    add_bottom_banner(slide,
                      "六项消融逐一回应审稿可能提出的方法/数据/物理/边界风险。")


def slide_18_writing_progress(prs):
    """论文进展页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COL_BG)
    draw_page_chrome(slide, 18, "WRITING", "论文进展",
                     "论文写作进展：v0.1 初稿已完成")
    add_subtitle(slide, "GPT 底稿 + Claude 组织 + 后整合双线修订",
                 top=Inches(1.50))

    # stage strip
    stages = [("v0.1", "主稿整合"),
              ("01", "作者确认"),
              ("02", "引用核验"),
              ("03", "图表定稿"),
              ("04", "全文压缩"),
              ("05", "模拟审稿"),
              ("06", "投稿材料")]
    top = Inches(2.10)
    h = Inches(0.95)
    total_w = Inches(12.4)
    n = len(stages)
    w_each = Inches((12.4 - 0.6) / n)
    gap = Inches(0.10)
    for i, (k, v) in enumerate(stages):
        l = MARGIN + (w_each + gap) * i
        color = COL_DEEP_GREEN if i == 0 else COL_MID_GREEN if i <= 1 else COL_TEAL
        bg = COL_LIGHT_GREEN if i == 0 else COL_SOFT_BG
        add_round_rect(slide, l, top, w_each, h, fill=bg, line=color, line_w=Pt(0.9), corner=0.06)
        add_text(slide, l, top + Inches(0.15), w_each, Inches(0.40),
                 [(k, {"size": 16, "bold": True, "color": color})],
                 align=PP_ALIGN.CENTER, vanchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, l, top + Inches(0.55), w_each, Inches(0.35),
                 [(v, {"size": 11, "color": COL_TEXT})],
                 align=PP_ALIGN.CENTER, vanchor=MSO_ANCHOR.MIDDLE)
        # arrow except last
        if i < n - 1:
            ax = l + w_each + Emu(2000)
            tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                          ax, top + Inches(0.30), gap, Inches(0.35))
            set_solid_fill(tri, COL_TEAL)

    # text under strip
    add_text(slide, MARGIN, top + h + Inches(0.18), Inches(12.4), Inches(0.4),
             [("六阶段后整合双线修订均已完成  →  当前进入 v0.2 前作者统一确认",
               {"size": 13, "bold": True, "color": COL_DEEP_GREEN})],
             align=PP_ALIGN.CENTER)

    # two cards
    top2 = Inches(4.20)
    add_round_rect(slide, MARGIN, top2, Inches(6.05), Inches(2.55),
                    fill=COL_SOFT_BG, line=COL_LIGHT_GREEN, line_w=Pt(1.0), corner=0.05)
    add_text(slide, MARGIN + Inches(0.20), top2 + Inches(0.10),
             Inches(5.7), Inches(0.4),
             [("当前状态", {"size": 14, "bold": True, "color": COL_DEEP_GREEN})])
    add_text(slide, MARGIN + Inches(0.20), top2 + Inches(0.55),
             Inches(5.7), Inches(2.0),
             [[("▍ 主稿 v0.1 已整合完成 (GPT 底稿 + Claude 组织)",
                {"size": 12, "color": COL_TEXT})],
              [("▍ 已不在从零写初稿阶段，转为定向修订",
                {"size": 12, "color": COL_TEXT})],
              [("▍ 暂不生成 v0.2：等 Blocking 项确认后再出新版本",
                {"size": 12, "color": COL_TEXT})]],
             line_spacing=1.45)

    x_right = MARGIN + Inches(6.35)
    add_round_rect(slide, x_right, top2, Inches(6.05), Inches(2.55),
                    fill=COL_ORANGE_BG, line=COL_ORANGE, line_w=Pt(1.0), corner=0.05)
    add_text(slide, x_right + Inches(0.20), top2 + Inches(0.10),
             Inches(5.7), Inches(0.4),
             [("待确认 Blocking 项", {"size": 14, "bold": True, "color": COL_ORANGE})])
    add_text(slide, x_right + Inches(0.20), top2 + Inches(0.55),
             Inches(5.7), Inches(2.0),
             [[("▍ 最终方法定义与关键数值口径",
                {"size": 12, "color": COL_TEXT})],
              [("▍ 引用核验、目标期刊",
                {"size": 12, "color": COL_TEXT})],
              [("▍ 数据 / 代码共享、作者与声明事实",
                {"size": 12, "color": COL_TEXT})],
              [("▍ 真实 ISAR / 观测数据：纳入主线 vs future work",
                {"size": 12, "color": COL_TEXT})]],
             line_spacing=1.45)

    add_bottom_banner(slide,
                      "汇报终点至 v0.1 整合稿；v0.2 之后的投稿材料细节本次不展开。")


# ---------- Build ----------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_2_significance(prs)
    slide_3_reliability(prs)
    slide_4_complementarity(prs)
    slide_5_value(prs)
    slide_6_physical_consistency(prs)
    slide_7_brdf_route(prs)
    slide_8_sampling_diag(prs)
    slide_9_three_end_closure(prs)
    slide_10_route_status(prs)
    slide_11_forward_validation(prs)
    slide_12_methods_definition(prs)
    slide_13_main_results(prs)
    slide_14_resnet_repositioning(prs)
    slide_15_robustness_image(prs)
    slide_16_robustness_ocs(prs)
    slide_17_ablation_six(prs)
    slide_18_writing_progress(prs)

    prs.save(str(OUT_PPTX))
    print(f"[OK] wrote {OUT_PPTX} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
